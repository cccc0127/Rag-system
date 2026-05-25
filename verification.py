"""Verification utility for the RAG privacy pipeline.

The script runs a small end-to-end pass over the knowledge base:

1. recursively loads documents and chunks them with privacy scoring;
2. embeds chunks, performs JL reduction, applies analytic Gaussian noise;
3. reports dimension changes and L2 norm ranges;
4. plots raw sensitivity score vs. calibrated sigma;
5. compares cosine similarities before and after noise for one query.
"""

from __future__ import annotations

import argparse
import os
from pathlib import Path
from typing import Dict, Iterable, List

import numpy as np

from chunk import iter_chunk_documents
from config import config
from dimension_reduction import JLProjector, l2_normalize
from gaussian_noise import AnalyticGaussianCalibrator, NoiseApplication, NoiseCalibration
from privacy_judge import PrivacyScorer


TEXT_LIKE_SUFFIXES = {"", ".txt", ".md", ".csv", ".log"}
SUPPORTED_BINARY_SUFFIXES = {".pdf", ".docx", ".xlsx", ".xls", ".pptx"}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Verify chunking, JL reduction, Gaussian noise, and retrieval loss."
    )
    parser.add_argument(
        "--knowledge-base",
        type=Path,
        default=config.REFERENCE_FOLDER,
        help="Knowledge-base directory to verify. Defaults to config.REFERENCE_FOLDER.",
    )
    parser.add_argument(
        "--embedding-model",
        default=config.EMBEDDING_MODEL,
        help="SentenceTransformer model path or name.",
    )
    parser.add_argument(
        "--query",
        default="What are the main topics discussed in these documents?",
        help="Query used for cosine-similarity stress testing.",
    )
    parser.add_argument("--max-chunks", type=int, default=100)
    parser.add_argument("--chunk-size", type=int, default=getattr(config, "CHUNK_SIZE", 1000))
    parser.add_argument("--overlap", type=int, default=getattr(config, "OVERLAP", 200))
    parser.add_argument("--jl-target-dim", type=int, default=getattr(config, "JL_TARGET_DIM", 256))
    parser.add_argument("--jl-epsilon", type=float, default=getattr(config, "JL_EPSILON", 0.3))
    parser.add_argument("--jl-seed", type=int, default=getattr(config, "JL_RANDOM_SEED", 42))
    parser.add_argument("--noise-seed", type=int, default=getattr(config, "DP_RANDOM_SEED", 42))
    parser.add_argument("--dp-delta", type=float, default=getattr(config, "DP_DELTA", 1e-5))
    parser.add_argument("--utility-scale", type=float, default=getattr(config, "DP_UTILITY_SCALE", 0.01))
    parser.add_argument(
        "--enable-nlp-privacy",
        action="store_true",
        help="Enable the optional zero-shot privacy classifier. Disabled by default for offline verification.",
    )
    parser.add_argument(
        "--plot-path",
        type=Path,
        default=Path("noise_distribution.png"),
        help="Where to save the raw_score/sigma trend plot.",
    )
    return parser.parse_args()


def iter_documents_recursive(root: Path) -> Iterable[Dict[str, str]]:
    """Yield supported files recursively, including extensionless Enron mail files."""
    if not root.exists():
        raise FileNotFoundError(f"Knowledge base does not exist: {root}")

    for path in sorted(p for p in root.rglob("*") if p.is_file()):
        if path.name.startswith("."):
            continue
        content = load_single_document(path)
        if content.strip():
            yield {"filename": str(path.relative_to(root)), "content": content}


def load_documents_recursive(root: Path) -> List[Dict[str, str]]:
    """Eager compatibility wrapper around iter_documents_recursive."""
    return list(iter_documents_recursive(root))


def count_iterated_documents(
    docs: Iterable[Dict[str, str]],
    counter: Dict[str, int],
) -> Iterable[Dict[str, str]]:
    for doc in docs:
        counter["count"] += 1
        yield doc


def load_single_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_LIKE_SUFFIXES or path.name.rstrip(".").isdigit():
        return read_text_lossy(path)
    if suffix == ".pdf":
        return read_pdf(path)
    if suffix == ".docx":
        return read_docx(path)
    if suffix in {".xlsx", ".xls"}:
        return read_excel(path)
    if suffix == ".pptx":
        return read_pptx(path)
    return ""


def read_text_lossy(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ModuleNotFoundError:
        print(f"Skipping PDF without pypdf installed: {path}")
        return ""

    try:
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        print(f"Skipping unreadable PDF {path}: {exc}")
        return ""


def read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ModuleNotFoundError:
        print(f"Skipping DOCX without python-docx installed: {path}")
        return ""

    try:
        doc = Document(str(path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    except Exception as exc:
        print(f"Skipping unreadable DOCX {path}: {exc}")
        return ""


def read_excel(path: Path) -> str:
    try:
        import pandas as pd
    except ModuleNotFoundError:
        print(f"Skipping spreadsheet without pandas installed: {path}")
        return ""

    try:
        return pd.read_excel(path).to_csv(index=False, sep="\t")
    except Exception as exc:
        print(f"Skipping unreadable spreadsheet {path}: {exc}")
        return ""


def read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ModuleNotFoundError:
        print(f"Skipping PPTX without python-pptx installed: {path}")
        return ""

    try:
        prs = Presentation(str(path))
        text: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
        return "\n".join(text)
    except Exception as exc:
        print(f"Skipping unreadable PPTX {path}: {exc}")
        return ""


def take_chunks(
    docs: Iterable[Dict[str, str]],
    max_chunks: int,
    chunk_size: int,
    overlap: int,
    enable_nlp_privacy: bool,
) -> List[Dict[str, object]]:
    scorer = PrivacyScorer(enable_nlp=enable_nlp_privacy)
    chunks: List[Dict[str, object]] = []
    for record in iter_chunk_documents(docs, chunk_size=chunk_size, overlap=overlap, scorer=scorer):
        chunks.append(record)
        if len(chunks) >= max_chunks:
            break
    return chunks


def apply_noise_with_calibrations(
    embeddings: np.ndarray,
    chunk_records: List[Dict[str, object]],
    calibrator: AnalyticGaussianCalibrator,
) -> tuple[np.ndarray, np.ndarray, np.ndarray, List[NoiseCalibration], np.ndarray]:
    applications: List[NoiseApplication] = [
        calibrator.apply_noise_with_diagnostics(
            vector,
            raw_score=float(record["raw_sensitivity_score"]),
        )
        for vector, record in zip(embeddings, chunk_records)
    ]
    noised = np.vstack([item.noised_vector for item in applications]).astype(np.float32)
    clipped = np.vstack([item.clipped_vector for item in applications]).astype(np.float32)
    noise = np.vstack([item.noise_vector for item in applications]).astype(np.float32)
    calibrations = [item.calibration for item in applications]
    sigma_per_dim = np.array([item.sigma_per_dim for item in applications], dtype=np.float32)
    return noised, clipped, noise, calibrations, sigma_per_dim


def plot_noise_distribution(
    calibrations: List[NoiseCalibration],
    plot_path: Path,
) -> None:
    try:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "matplotlib is required to save noise_distribution.png. "
            "Install project requirements with: pip install -r requirements.txt"
        ) from exc

    raw_scores = np.array([item.raw_score for item in calibrations], dtype=np.float32)
    sigmas = np.array([item.sigma for item in calibrations], dtype=np.float32)
    order = np.argsort(raw_scores)

    plt.figure(figsize=(9, 5))
    plt.scatter(raw_scores, sigmas, alpha=0.72, label="chunk")
    plt.plot(raw_scores[order], sigmas[order], color="#d95f02", linewidth=2, label="sorted trend")
    plt.xlabel("raw_score")
    plt.ylabel("sigma")
    plt.title("Analytic Gaussian Noise Calibration")
    plt.grid(True, alpha=0.25)
    plt.legend()
    plt.tight_layout()
    plot_path.parent.mkdir(parents=True, exist_ok=True)
    plt.savefig(plot_path, dpi=160)
    plt.close()


def cosine_scores(query_vector: np.ndarray, matrix: np.ndarray) -> np.ndarray:
    query = l2_normalize(query_vector)[0]
    docs = l2_normalize(matrix)
    return docs @ query


def paired_cosine(left: np.ndarray, right: np.ndarray) -> np.ndarray:
    left_norm = l2_normalize(left)
    right_norm = l2_normalize(right)
    return np.sum(left_norm * right_norm, axis=1)


def l2_norms(vectors: np.ndarray) -> np.ndarray:
    return np.linalg.norm(vectors, ord=2, axis=1)


def format_min_mean_max(values: np.ndarray) -> str:
    return f"{values.min():.6f} / {values.mean():.6f} / {values.max():.6f}"


def print_array_stats(name: str, vectors: np.ndarray) -> None:
    print(f"{name} shape:".ljust(40), vectors.shape)
    print(f"{name} L2 min/mean/max:".ljust(40), format_min_mean_max(l2_norms(vectors)))
    print(f"{name} value mean/std:".ljust(40), f"{vectors.mean():.6f} / {vectors.std():.6f}")


def print_dimension_report(
    raw_embeddings: np.ndarray,
    reduced_raw_embeddings: np.ndarray,
    clipped_embeddings: np.ndarray,
    noised_embeddings: np.ndarray,
    final_noised_embeddings: np.ndarray,
) -> None:
    print("\n=== Vector dimension changes ===")
    print(f"Chunk count:                         {raw_embeddings.shape[0]}")
    print(f"Raw embedding matrix:                {raw_embeddings.shape}")
    print(f"After JL reduction (raw baseline):   {reduced_raw_embeddings.shape}")
    print(f"After L2 clipping:                   {clipped_embeddings.shape}")
    print(f"After Gaussian noise:                {noised_embeddings.shape}")
    print(f"After final normalization:           {final_noised_embeddings.shape}")


def print_norm_report(
    raw_embeddings: np.ndarray,
    reduced_raw_embeddings: np.ndarray,
    clipped_embeddings: np.ndarray,
    noise_vectors: np.ndarray,
    noised_embeddings: np.ndarray,
    final_noised_embeddings: np.ndarray,
) -> None:
    noised_norms = np.linalg.norm(noised_embeddings, ord=2, axis=1)
    final_norms = np.linalg.norm(final_noised_embeddings, ord=2, axis=1)
    signal_norms = l2_norms(clipped_embeddings)
    noise_norms = l2_norms(noise_vectors)
    noise_signal_ratio = noise_norms / np.maximum(signal_norms, 1e-12)
    noised_in_unit_ball = bool(np.all((noised_norms >= 0.0) & (noised_norms <= 1.0)))
    final_in_unit_ball = bool(np.all((final_norms >= 0.0) & (final_norms <= 1.0 + 1e-6)))

    print("\n=== Vector statistics ===")
    print_array_stats("Raw embedding", raw_embeddings)
    print_array_stats("JL raw embedding", reduced_raw_embeddings)
    print_array_stats("Clipped signal", clipped_embeddings)
    print_array_stats("Gaussian noise", noise_vectors)
    print_array_stats("Noised pre-normalize", noised_embeddings)
    print_array_stats("Final noised", final_noised_embeddings)

    print("\n=== L2 norm checks ===")
    print(f"Clipped signal L2 min/mean/max:      {format_min_mean_max(signal_norms)}")
    print(f"Pure noise L2 min/mean/max:          {format_min_mean_max(noise_norms)}")
    print(f"Noise/Signal ratio min/mean/max:     {format_min_mean_max(noise_signal_ratio)}")
    print(f"Noised vector L2 min/mean/max:       {format_min_mean_max(noised_norms)}")
    print(f"Noised vectors all in [0, 1]:        {noised_in_unit_ball}")
    print(f"Final noised L2 min/mean/max:        {format_min_mean_max(final_norms)}")
    print(f"Final noised vectors all in [0, 1]:  {final_in_unit_ball}")
    if not noised_in_unit_ball:
        print("Note: Gaussian noise is unbounded, so post-noise vectors are not guaranteed to stay in the unit ball.")


def print_calibration_report(calibrations: List[NoiseCalibration], sigma_per_dim: np.ndarray) -> None:
    raw_scores = np.array([item.raw_score for item in calibrations], dtype=np.float32)
    epsilons = np.array([item.epsilon for item in calibrations], dtype=np.float32)
    local_sensitivities = np.array([item.local_sensitivity for item in calibrations], dtype=np.float32)
    sigmas = np.array([item.sigma for item in calibrations], dtype=np.float32)
    solved = sum(item.solved_by_brentq for item in calibrations)

    print("\n=== raw_score -> sigma calibration ===")
    print(f"Recorded pairs:                      {len(calibrations)}")
    print(f"raw_score min/mean/max:              {raw_scores.min():.6f} / {raw_scores.mean():.6f} / {raw_scores.max():.6f}")
    print(f"epsilon min/mean/max:                {format_min_mean_max(epsilons)}")
    print(f"local_sensitivity min/mean/max:      {format_min_mean_max(local_sensitivities)}")
    print(f"sigma min/mean/max:                  {sigmas.min():.6f} / {sigmas.mean():.6f} / {sigmas.max():.6f}")
    print(f"sigma_per_dim min/mean/max:          {format_min_mean_max(sigma_per_dim)}")
    print(f"Analytic root solved by brentq:       {solved}/{len(calibrations)}")
    print("First 10 pairs:")
    for idx, item in enumerate(calibrations[:10], start=1):
        print(
            f"  {idx:02d}. raw_score={item.raw_score:.6f}, "
            f"epsilon={item.epsilon:.6f}, sigma={item.sigma:.6f}, "
            f"sigma_per_dim={sigma_per_dim[idx - 1]:.8f}"
        )


def print_similarity_report(
    query: str,
    raw_query: np.ndarray,
    query_reduced: np.ndarray,
    raw_embeddings: np.ndarray,
    reduced_raw_embeddings: np.ndarray,
    clipped_embeddings: np.ndarray,
    noised_embeddings: np.ndarray,
    final_noised_embeddings: np.ndarray,
    chunk_records: List[Dict[str, object]],
    top_k: int = 5,
) -> None:
    raw_1024_scores = cosine_scores(raw_query, raw_embeddings)
    reduced_scores = cosine_scores(query_reduced, reduced_raw_embeddings)
    noised_scores = cosine_scores(query_reduced, final_noised_embeddings)
    jl_abs_diffs = np.abs(reduced_scores - raw_1024_scores)
    raw_scores = reduced_scores
    diffs = noised_scores - raw_scores
    abs_diffs = np.abs(diffs)
    dp_pre_norm_cosine = paired_cosine(clipped_embeddings, noised_embeddings)
    dp_final_cosine = paired_cosine(reduced_raw_embeddings, final_noised_embeddings)

    print("\n=== Retrieval stress test: cosine similarity loss ===")
    print(f"Query:                               {query}")
    print(f"Mean raw 1024d cosine:               {raw_1024_scores.mean():.6f}")
    print(f"Mean JL raw 256d cosine:             {reduced_scores.mean():.6f}")
    print(f"Mean final noised cosine:            {noised_scores.mean():.6f}")
    print(f"JL query-score abs drift mean/max:   {jl_abs_diffs.mean():.6f} / {jl_abs_diffs.max():.6f}")
    print(f"DP score abs drift mean/max:         {abs_diffs.mean():.6f} / {abs_diffs.max():.6f}")
    print(f"DP score Pearson correlation:        {safe_corrcoef(raw_scores, noised_scores):.6f}")
    print(f"Pre-normalize signal/noised cosine:  {format_min_mean_max(dp_pre_norm_cosine)}")
    print(f"Final reduced/noised cosine:         {format_min_mean_max(dp_final_cosine)}")

    print(f"\nTop {top_k} raw-vs-noised score comparison:")
    top_indices = np.argsort(raw_scores)[::-1][:top_k]
    for rank, idx in enumerate(top_indices, start=1):
        record = chunk_records[int(idx)]
        name = f"{record['filename']}#chunk-{record['chunk_id']}"
        print(
            f"  {rank}. {name}: raw={raw_scores[idx]:.6f}, "
            f"noised={noised_scores[idx]:.6f}, diff={diffs[idx]:+.6f}"
        )


def safe_corrcoef(left: np.ndarray, right: np.ndarray) -> float:
    if left.size < 2 or np.std(left) == 0.0 or np.std(right) == 0.0:
        return float("nan")
    return float(np.corrcoef(left, right)[0, 1])


def load_embedding_model(model_path: str):
    try:
        from sentence_transformers import SentenceTransformer
    except Exception as exc:
        raise RuntimeError(
            "Failed to import sentence-transformers. This is usually caused by "
            "missing packages or incompatible dependency versions. Try: "
            "python3 -m pip install --upgrade -r requirements.txt"
        ) from exc

    return SentenceTransformer(str(model_path), device=config.EMBEDDING_DEVICE)


def main() -> None:
    args = parse_args()
    os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")

    print("=== Verification setup ===")
    print(f"Knowledge base:                      {args.knowledge_base}")
    print(f"Embedding model:                     {args.embedding_model}")
    print(f"Max chunks:                          {args.max_chunks}")

    document_counter = {"count": 0}
    docs = count_iterated_documents(
        iter_documents_recursive(args.knowledge_base),
        document_counter,
    )

    chunk_records = take_chunks(
        docs,
        max_chunks=args.max_chunks,
        chunk_size=args.chunk_size,
        overlap=args.overlap,
        enable_nlp_privacy=args.enable_nlp_privacy,
    )
    if not chunk_records:
        raise RuntimeError("No chunks were produced from the loaded documents.")

    print(f"Scanned readable documents:          {document_counter['count']}")
    print(f"Processed chunks:                    {len(chunk_records)}")

    texts = [str(record["content"]) for record in chunk_records]
    embedding_model = load_embedding_model(str(args.embedding_model))
    raw_embeddings = np.asarray(
        embedding_model.encode(texts, batch_size=16, show_progress_bar=True),
        dtype=np.float32,
    )

    projector = JLProjector(
        target_dim=args.jl_target_dim,
        eps=args.jl_epsilon,
        random_state=args.jl_seed,
    )
    reduced_raw_embeddings = projector.fit_transform(raw_embeddings)
    calibrator = AnalyticGaussianCalibrator(
        delta=args.dp_delta,
        utility_scale=args.utility_scale,
        random_state=args.noise_seed,
    )
    (
        noised_embeddings,
        clipped_embeddings,
        noise_vectors,
        calibrations,
        sigma_per_dim,
    ) = apply_noise_with_calibrations(
        reduced_raw_embeddings,
        chunk_records,
        calibrator,
    )
    final_noised_embeddings = l2_normalize(noised_embeddings)

    query_embedding = np.asarray(embedding_model.encode([args.query]), dtype=np.float32)
    query_reduced = projector.transform(query_embedding)

    print_dimension_report(
        raw_embeddings,
        reduced_raw_embeddings,
        clipped_embeddings,
        noised_embeddings,
        final_noised_embeddings,
    )
    print_norm_report(
        raw_embeddings,
        reduced_raw_embeddings,
        clipped_embeddings,
        noise_vectors,
        noised_embeddings,
        final_noised_embeddings,
    )
    print_calibration_report(calibrations, sigma_per_dim)
    plot_noise_distribution(calibrations, args.plot_path)
    print(f"\nSaved noise trend plot:              {args.plot_path}")
    print_similarity_report(
        args.query,
        query_embedding,
        query_reduced,
        raw_embeddings,
        reduced_raw_embeddings,
        clipped_embeddings,
        noised_embeddings,
        final_noised_embeddings,
        chunk_records,
    )


if __name__ == "__main__":
    main()
