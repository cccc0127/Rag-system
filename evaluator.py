"""End-to-end utility evaluation for the DP-RAG retrieval pipeline.

This script samples chunks, builds raw/JL/noised embeddings in memory, and
compares Top-5 retrieval behavior before and after DP noise.
"""

from __future__ import annotations

import argparse
import os
import re
import textwrap
import time
from pathlib import Path
from typing import Dict, Iterable, List, Sequence

os.environ.setdefault("MPLCONFIGDIR", "/tmp/matplotlib")

import matplotlib.pyplot as plt
import numpy as np

from chunk import iter_chunk_documents
from config import config
from dimension_reduction import JLProjector, l2_normalize
from gaussian_noise import AnalyticGaussianCalibrator, NoiseApplication
from privacy_judge import PrivacyScorer


TEXT_LIKE_SUFFIXES = {"", ".txt", ".md", ".csv", ".log"}
WORD_RE = re.compile(r"[A-Za-z][A-Za-z0-9_'-]{2,}|[\u4e00-\u9fff]{2,}")
UTILITY_SCALE_LIST = [0.001, 0.01, 0.05, 0.1, 0.2, 0.5]
RESULT_DIR = Path(__file__).resolve().parent / "Result_picture"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Evaluate DP-RAG retrieval utility on sampled knowledge-base chunks."
    )
    parser.add_argument("--knowledge-base", type=Path, default=config.REFERENCE_FOLDER)
    parser.add_argument("--embedding-model", default=config.EMBEDDING_MODEL)
    parser.add_argument("--sample-chunks", type=int, default=100)
    parser.add_argument("--num-queries", type=int, default=5)
    parser.add_argument("--top-k", type=int, default=5)
    parser.add_argument("--chunk-size", type=int, default=getattr(config, "CHUNK_SIZE", 1000))
    parser.add_argument("--overlap", type=int, default=getattr(config, "OVERLAP", 200))
    parser.add_argument("--jl-target-dim", type=int, default=getattr(config, "JL_TARGET_DIM", 256))
    parser.add_argument("--jl-epsilon", type=float, default=getattr(config, "JL_EPSILON", 0.3))
    parser.add_argument("--jl-seed", type=int, default=getattr(config, "JL_RANDOM_SEED", 42))
    parser.add_argument("--noise-seed", type=int, default=getattr(config, "DP_RANDOM_SEED", 42))
    parser.add_argument("--dp-delta", type=float, default=getattr(config, "DP_DELTA", 1e-5))
    parser.add_argument("--utility-scale", type=float, default=getattr(config, "DP_UTILITY_SCALE", 0.01))
    parser.add_argument("--query-seed", type=int, default=2026)
    parser.add_argument("--visual-queries", type=int, default=1)
    parser.add_argument("--visual-text-chars", type=int, default=400)
    parser.add_argument(
        "--enable-nlp-privacy",
        action="store_true",
        help="Enable optional zero-shot privacy scoring. Disabled by default for offline evaluation.",
    )
    return parser.parse_args()


def iter_documents_recursive(root: Path) -> Iterable[Dict[str, str]]:
    if not root.exists():
        raise FileNotFoundError(f"Knowledge base does not exist: {root}")

    for path in sorted(p for p in root.rglob("*") if p.is_file()):
        if path.name.startswith("."):
            continue
        content = load_single_document(path)
        if content.strip():
            yield {"filename": str(path.relative_to(root)), "content": content}


def load_single_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in TEXT_LIKE_SUFFIXES or path.name.rstrip(".").isdigit():
        return read_text_lossy(path)
    if suffix == ".pdf":
        return read_pdf(path)
    if suffix == ".docx":
        return read_docx(path)
    if suffix in {".xlsx", ".xls"}:
        return read_excel(path)
    if suffix == ".pptx":
        return read_pptx(path)
    return ""


def read_text_lossy(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def read_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ModuleNotFoundError:
        return ""
    try:
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return ""


def read_docx(path: Path) -> str:
    try:
        from docx import Document
    except ModuleNotFoundError:
        return ""
    try:
        doc = Document(str(path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    except Exception:
        return ""


def read_excel(path: Path) -> str:
    try:
        import pandas as pd
    except ModuleNotFoundError:
        return ""
    try:
        return pd.read_excel(path).to_csv(index=False, sep="\t")
    except Exception:
        return ""


def read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ModuleNotFoundError:
        return ""
    try:
        prs = Presentation(str(path))
        text: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.extend(paragraph.text for paragraph in shape.text_frame.paragraphs)
        return "\n".join(text)
    except Exception:
        return ""


def count_iterated_documents(
    docs: Iterable[Dict[str, str]],
    counter: Dict[str, int],
) -> Iterable[Dict[str, str]]:
    for doc in docs:
        counter["count"] += 1
        yield doc


def sample_chunks(
    docs: Iterable[Dict[str, str]],
    max_chunks: int,
    chunk_size: int,
    overlap: int,
    enable_nlp_privacy: bool,
) -> List[Dict[str, object]]:
    scorer = PrivacyScorer(enable_nlp=enable_nlp_privacy)
    chunks: List[Dict[str, object]] = []
    for record in iter_chunk_documents(docs, chunk_size=chunk_size, overlap=overlap, scorer=scorer):
        chunks.append(record)
        if len(chunks) >= max_chunks:
            break
    return chunks


def load_embedding_model(model_path: str):
    try:
        from sentence_transformers import SentenceTransformer
    except Exception as exc:
        raise RuntimeError(
            "Failed to import sentence-transformers. Try: "
            "python3 -m pip install --upgrade -r requirements.txt"
        ) from exc
    return SentenceTransformer(str(model_path), device=config.EMBEDDING_DEVICE)


def apply_dp_noise(
    reduced_embeddings: np.ndarray,
    chunk_records: Sequence[Dict[str, object]],
    calibrator: AnalyticGaussianCalibrator,
) -> tuple[np.ndarray, np.ndarray, np.ndarray, np.ndarray, np.ndarray, np.ndarray]:
    applications: List[NoiseApplication] = [
        calibrator.apply_noise_with_diagnostics(
            vector,
            raw_score=float(record["raw_sensitivity_score"]),
        )
        for vector, record in zip(reduced_embeddings, chunk_records)
    ]
    clipped = np.vstack([item.clipped_vector for item in applications]).astype(np.float32)
    noise = np.vstack([item.noise_vector for item in applications]).astype(np.float32)
    noised = np.vstack([item.noised_vector for item in applications]).astype(np.float32)
    final_noised = l2_normalize(noised)
    sigmas = np.array([item.calibration.sigma for item in applications], dtype=np.float32)
    sigma_per_dim = np.array([item.sigma_per_dim for item in applications], dtype=np.float32)
    epsilons = np.array([item.calibration.epsilon for item in applications], dtype=np.float32)
    return final_noised, clipped, noise, sigmas, sigma_per_dim, epsilons


def generate_random_queries(
    chunk_records: Sequence[Dict[str, object]],
    num_queries: int,
    seed: int,
) -> List[str]:
    rng = np.random.default_rng(seed)
    if not chunk_records:
        return []

    queries: List[str] = []
    selected = rng.choice(len(chunk_records), size=min(num_queries, len(chunk_records)), replace=False)
    for idx in selected:
        content = str(chunk_records[int(idx)]["content"])
        terms = [term.lower() for term in WORD_RE.findall(content)]
        unique_terms = list(dict.fromkeys(term for term in terms if len(term.strip()) >= 3))
        if len(unique_terms) >= 4:
            take = min(8, len(unique_terms))
            query_terms = rng.choice(unique_terms, size=take, replace=False)
            queries.append(" ".join(query_terms))
        else:
            snippet = " ".join(content.split()[:16]).strip()
            queries.append(snippet or "general document topic")

    while len(queries) < num_queries:
        queries.append(f"general document topic {len(queries) + 1}")
    return queries[:num_queries]


def cosine_scores(query_vector: np.ndarray, matrix: np.ndarray) -> np.ndarray:
    query = l2_normalize(query_vector)[0]
    docs = l2_normalize(matrix)
    return docs @ query


def paired_cosine(left: np.ndarray, right: np.ndarray) -> np.ndarray:
    left_norm = l2_normalize(left)
    right_norm = l2_normalize(right)
    return np.sum(left_norm * right_norm, axis=1)


def top_k_indices(scores: np.ndarray, top_k: int) -> np.ndarray:
    top_k = min(top_k, scores.size)
    return np.argsort(scores)[::-1][:top_k]


def overlap_at_k(raw_top: np.ndarray, noised_top: np.ndarray, k: int) -> float:
    raw_set = set(int(idx) for idx in raw_top[:k])
    noised_set = set(int(idx) for idx in noised_top[:k])
    return len(raw_set & noised_set) / max(1, min(k, len(raw_set)))


def reciprocal_rank_at_k(target_id: int, ranked_ids: np.ndarray, k: int) -> float:
    for rank, idx in enumerate(ranked_ids[:k], start=1):
        if int(idx) == int(target_id):
            return 1.0 / rank
    return 0.0


def safe_pearson(left: np.ndarray, right: np.ndarray) -> float:
    if left.size < 2 or np.std(left) == 0.0 or np.std(right) == 0.0:
        return float("nan")
    return float(np.corrcoef(left, right)[0, 1])


def l2_norms(vectors: np.ndarray) -> np.ndarray:
    return np.linalg.norm(vectors, ord=2, axis=1)


def truncate(text: str, width: int) -> str:
    text = " ".join(text.split())
    if len(text) <= width:
        return text
    return text[: max(0, width - 3)] + "..."


def print_table(headers: Sequence[str], rows: Sequence[Sequence[object]]) -> None:
    table = [[str(value) for value in row] for row in rows]
    widths = [
        max(len(str(header)), *(len(row[idx]) for row in table)) if table else len(str(header))
        for idx, header in enumerate(headers)
    ]
    border = "+-" + "-+-".join("-" * width for width in widths) + "-+"
    print(border)
    print("| " + " | ".join(str(header).ljust(widths[idx]) for idx, header in enumerate(headers)) + " |")
    print(border)
    for row in table:
        print("| " + " | ".join(row[idx].ljust(widths[idx]) for idx in range(len(headers))) + " |")
    print(border)


def clean_for_panel(text: str, max_chars: int) -> str:
    text = "\n".join(line.rstrip() for line in str(text).strip().splitlines())
    if len(text) <= max_chars:
        return text
    return text[: max(0, max_chars - 3)].rstrip() + "..."


def wrap_panel_field(label: str, value: str, width: int) -> List[str]:
    prefix = f"{label}: "
    available = max(20, width - len(prefix))
    lines: List[str] = []
    paragraphs = str(value).splitlines() or [""]

    for paragraph_idx, paragraph in enumerate(paragraphs):
        wrapped = textwrap.wrap(
            paragraph,
            width=available if not lines else width,
            replace_whitespace=False,
            drop_whitespace=False,
        ) or [""]
        for line_idx, line in enumerate(wrapped):
            if not lines and paragraph_idx == 0 and line_idx == 0:
                lines.append(f"{prefix}{line}")
            elif line_idx == 0 and paragraph_idx > 0:
                lines.append("")
                lines.append(line)
            else:
                lines.append((" " * len(prefix) if len(lines) == 1 else "") + line)
    return lines


def print_visual_alignment_panel(
    query_id: int,
    query: str,
    top1_id: int,
    raw_text: str,
    max_chars: int,
    width: int = 96,
) -> None:
    title = f"Visual Semantic Alignment Verification | Query {query_id}"
    inner_width = width - 4
    border = "+" + "-" * (width - 2) + "+"
    print(border)
    print("| " + title.center(inner_width) + " |")
    print(border)

    fields = [
        ("Query", query),
        ("Raw Top-1 Chunk ID", str(top1_id)),
        ("Raw Text", clean_for_panel(raw_text, max_chars)),
    ]
    for field_idx, (label, value) in enumerate(fields):
        for line in wrap_panel_field(label, value, inner_width):
            print("| " + line[:inner_width].ljust(inner_width) + " |")
        if field_idx < len(fields) - 1:
            print("| " + "-" * inner_width + " |")
    print(border)


def prepare_evaluation_context(args: argparse.Namespace) -> Dict[str, object]:
    os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")

    doc_counter = {"count": 0}
    docs = count_iterated_documents(iter_documents_recursive(args.knowledge_base), doc_counter)
    chunk_records = sample_chunks(
        docs,
        max_chunks=args.sample_chunks,
        chunk_size=args.chunk_size,
        overlap=args.overlap,
        enable_nlp_privacy=args.enable_nlp_privacy,
    )
    if not chunk_records:
        raise RuntimeError("No chunks were sampled from the knowledge base.")

    texts = [str(record["content"]) for record in chunk_records]
    embedding_model = load_embedding_model(str(args.embedding_model))
    raw_embeddings = np.asarray(
        embedding_model.encode(texts, batch_size=16, show_progress_bar=True),
        dtype=np.float32,
    )

    projector = JLProjector(
        target_dim=args.jl_target_dim,
        eps=args.jl_epsilon,
        random_state=args.jl_seed,
    )
    reduced_raw_embeddings = projector.fit_transform(raw_embeddings)

    queries = generate_random_queries(chunk_records, args.num_queries, args.query_seed)
    query_embeddings = np.asarray(
        embedding_model.encode(queries, batch_size=args.num_queries, show_progress_bar=False),
        dtype=np.float32,
    )
    query_reduced = projector.transform(query_embeddings)

    return {
        "doc_counter": doc_counter,
        "chunk_records": chunk_records,
        "raw_embeddings": raw_embeddings,
        "reduced_raw_embeddings": reduced_raw_embeddings,
        "queries": queries,
        "query_reduced": query_reduced,
    }


def evaluate_utility_scale(
    args: argparse.Namespace,
    context: Dict[str, object],
    utility_scale: float,
) -> Dict[str, float]:
    config.DP_UTILITY_SCALE = float(utility_scale)

    doc_counter = context["doc_counter"]
    chunk_records = context["chunk_records"]
    raw_embeddings = context["raw_embeddings"]
    reduced_raw_embeddings = context["reduced_raw_embeddings"]
    queries = context["queries"]
    query_reduced = context["query_reduced"]

    assert isinstance(doc_counter, dict)
    assert isinstance(chunk_records, list)
    assert isinstance(raw_embeddings, np.ndarray)
    assert isinstance(reduced_raw_embeddings, np.ndarray)
    assert isinstance(queries, list)
    assert isinstance(query_reduced, np.ndarray)

    calibrator = AnalyticGaussianCalibrator(
        delta=args.dp_delta,
        utility_scale=float(config.DP_UTILITY_SCALE),
        random_state=args.noise_seed,
    )
    noise_start = time.perf_counter()
    (
        final_noised_embeddings,
        clipped_embeddings,
        noise_vectors,
        sigmas,
        sigma_per_dim,
        epsilons,
    ) = apply_dp_noise(
        reduced_raw_embeddings,
        chunk_records,
        calibrator,
    )
    dp_noise_time = time.perf_counter() - noise_start

    signal_norms = l2_norms(clipped_embeddings)
    noise_norms = l2_norms(noise_vectors)
    noise_signal_ratios = noise_norms / np.maximum(signal_norms, 1e-12)
    mean_noise_signal_ratio = float(np.mean(noise_signal_ratios))
    direction_cosines = paired_cosine(reduced_raw_embeddings, final_noised_embeddings)

    rows = []
    overlap_values: Dict[int, List[float]] = {1: [], 3: [], 5: [], 10: []}
    pearson_values: List[float] = []
    drift_values: List[float] = []
    max_drift_values: List[float] = []
    mrr5_values: List[float] = []
    raw_retrieval_times: List[float] = []
    noised_retrieval_times: List[float] = []
    visual_samples: List[Dict[str, object]] = []
    retrieval_depth = max(args.top_k, 10, 5)

    for query_id, (query, query_vec) in enumerate(zip(queries, query_reduced), start=1):
        raw_start = time.perf_counter()
        raw_scores = cosine_scores(query_vec, reduced_raw_embeddings)
        raw_top_full = top_k_indices(raw_scores, retrieval_depth)
        raw_retrieval_times.append(time.perf_counter() - raw_start)

        noised_start = time.perf_counter()
        noised_scores = cosine_scores(query_vec, final_noised_embeddings)
        noised_top_full = top_k_indices(noised_scores, retrieval_depth)
        noised_retrieval_times.append(time.perf_counter() - noised_start)

        raw_top = raw_top_full[: min(args.top_k, raw_top_full.size)]
        noised_top = noised_top_full[: min(args.top_k, noised_top_full.size)]
        if raw_top.size > 0 and query_id <= args.visual_queries:
            top1_id = int(raw_top[0])
            visual_samples.append(
                {
                    "query_id": query_id,
                    "query": query,
                    "top1_id": top1_id,
                    "raw_text": str(chunk_records[top1_id]["content"]),
                }
            )
        for k in overlap_values:
            overlap_values[k].append(overlap_at_k(raw_top_full, noised_top_full, k))
        overlap = overlap_at_k(raw_top_full, noised_top_full, args.top_k)
        mrr5_values.append(reciprocal_rank_at_k(int(raw_top_full[0]), noised_top_full, 5))
        pearson = safe_pearson(raw_scores, noised_scores)
        abs_drift = np.abs(noised_scores - raw_scores)
        drift = float(np.mean(abs_drift))
        max_drift = float(np.max(abs_drift))

        pearson_values.append(float(pearson))
        drift_values.append(drift)
        max_drift_values.append(max_drift)
        rows.append(
            [
                query_id,
                truncate(query, 38),
                f"{overlap:.3f}",
                f"{pearson:.6f}",
                f"{drift:.6f}",
                f"{reciprocal_rank_at_k(int(raw_top_full[0]), noised_top_full, 5):.3f}",
                ",".join(str(idx) for idx in raw_top.tolist()),
                ",".join(str(idx) for idx in noised_top.tolist()),
            ]
        )

    print("\nDP-RAG Retrieval Utility Evaluation")
    print("=" * 78)
    print(f"Utility scale:              {utility_scale}")
    print(f"Knowledge base:             {args.knowledge_base}")
    print(f"Scanned readable documents: {doc_counter['count']}")
    print(f"Sampled chunks:             {len(chunk_records)}")
    print(f"Raw embedding shape:        {raw_embeddings.shape}")
    print(f"JL raw embedding shape:     {reduced_raw_embeddings.shape}")
    print(f"Final noised shape:         {final_noised_embeddings.shape}")
    print(f"Queries:                    {len(queries)}")
    print(f"Top-K:                      {args.top_k}")

    print("\nPer-query Summary Report")
    print_table(
        ["Q", "Query", f"Overlap@{args.top_k}", "Pearson", "MeanAbsDrift", "MRR@5", "R_raw", "R_noised"],
        rows,
    )

    mean_overlap1 = float(np.mean(overlap_values[1]))
    mean_overlap3 = float(np.mean(overlap_values[3]))
    mean_overlap5 = float(np.mean(overlap_values[5]))
    mean_overlap10 = float(np.mean(overlap_values[10]))
    mean_overlap = mean_overlap5
    mean_pearson = float(np.nanmean(pearson_values))
    mean_drift = float(np.mean(drift_values))
    max_abs_drift = float(np.max(max_drift_values))
    mean_mrr5 = float(np.mean(mrr5_values))
    mean_direction_cosine = float(np.mean(direction_cosines))
    min_direction_cosine = float(np.min(direction_cosines))
    max_direction_cosine = float(np.max(direction_cosines))
    mean_raw_retrieval_time = float(np.mean(raw_retrieval_times))
    mean_noised_retrieval_time = float(np.mean(noised_retrieval_times))

    print("\nAggregate Metrics")
    print_table(
        ["Metric", "Value"],
        [
            ["Mean Noise/Signal Ratio", f"{mean_noise_signal_ratio:.6f}"],
            ["Mean Overlap@1", f"{mean_overlap1:.6f}"],
            ["Mean Overlap@3", f"{mean_overlap3:.6f}"],
            ["Mean Overlap@5", f"{mean_overlap5:.6f}"],
            ["Mean Overlap@10", f"{mean_overlap10:.6f}"],
            ["Mean MRR@5", f"{mean_mrr5:.6f}"],
            ["Mean Pearson Correlation", f"{mean_pearson:.6f}"],
            ["Mean Absolute Drift", f"{mean_drift:.6f}"],
            ["Max Absolute Drift", f"{max_abs_drift:.6f}"],
            ["Direction Cosine Min/Mean/Max", f"{min_direction_cosine:.6f} / {mean_direction_cosine:.6f} / {max_direction_cosine:.6f}"],
            ["Mean Sigma", f"{float(np.mean(sigmas)):.6f}"],
            ["Max Sigma", f"{float(np.max(sigmas)):.6f}"],
            ["Mean Sigma_per_dim", f"{float(np.mean(sigma_per_dim)):.8f}"],
            ["Mean Epsilon", f"{float(np.mean(epsilons)):.6f}"],
            ["Min Epsilon", f"{float(np.min(epsilons)):.6f}"],
            ["DP Noise Injection Time", f"{dp_noise_time:.6f}s"],
            ["Mean Raw Retrieval Time / Query", f"{mean_raw_retrieval_time:.6f}s"],
            ["Mean Noised Retrieval Time / Query", f"{mean_noised_retrieval_time:.6f}s"],
            ["Noise L2 Min/Mean/Max", f"{noise_norms.min():.6f} / {noise_norms.mean():.6f} / {noise_norms.max():.6f}"],
            ["Signal L2 Min/Mean/Max", f"{signal_norms.min():.6f} / {signal_norms.mean():.6f} / {signal_norms.max():.6f}"],
        ],
    )

    if visual_samples:
        print("\nVisual Semantic Alignment Verification")
        for sample in visual_samples:
            print_visual_alignment_panel(
                query_id=int(sample["query_id"]),
                query=str(sample["query"]),
                top1_id=int(sample["top1_id"]),
                raw_text=str(sample["raw_text"]),
                max_chars=args.visual_text_chars,
            )

    return {
        "utility_scale": float(utility_scale),
        "mean_nsr": mean_noise_signal_ratio,
        "mean_overlap": mean_overlap,
        "mean_overlap1": mean_overlap1,
        "mean_overlap3": mean_overlap3,
        "mean_overlap5": mean_overlap5,
        "mean_overlap10": mean_overlap10,
        "mean_pearson": mean_pearson,
        "mean_drift": mean_drift,
        "max_drift": max_abs_drift,
        "mean_mrr5": mean_mrr5,
        "mean_direction_cosine": mean_direction_cosine,
        "dp_noise_time": dp_noise_time,
        "mean_raw_retrieval_time": mean_raw_retrieval_time,
        "mean_noised_retrieval_time": mean_noised_retrieval_time,
    }


def plot_privacy_utility_tradeoff(results: Sequence[Dict[str, float]], output_path: Path) -> None:
    utility_scales = np.array([item["utility_scale"] for item in results], dtype=np.float64)
    mean_nsr = np.maximum(
        np.array([item["mean_nsr"] for item in results], dtype=np.float64),
        1e-12,
    )
    mean_overlap = np.array([item["mean_overlap"] for item in results], dtype=np.float64)
    mean_pearson = np.array([item["mean_pearson"] for item in results], dtype=np.float64)

    fig, ax_left = plt.subplots(figsize=(9.5, 5.8), dpi=160)
    ax_right = ax_left.twinx()

    line_overlap, = ax_left.plot(
        utility_scales,
        mean_overlap,
        marker="o",
        linewidth=2.2,
        color="#2563eb",
        label="Mean Overlap@5",
    )
    line_pearson, = ax_left.plot(
        utility_scales,
        mean_pearson,
        marker="s",
        linewidth=2.2,
        color="#16a34a",
        label="Mean Pearson Correlation",
    )
    line_nsr, = ax_right.plot(
        utility_scales,
        mean_nsr,
        marker="^",
        linewidth=2.2,
        color="#f97316",
        label="Mean Noise/Signal Ratio",
    )

    ax_left.set_title("Privacy-Utility Tradeoff under Dynamic DP Noise", pad=12)
    ax_left.set_xlabel("Utility Scale")
    ax_left.set_ylabel("Retrieval Utility: Overlap@5 / Pearson")
    ax_right.set_ylabel("Mean Noise/Signal Ratio (log scale)")
    ax_left.set_ylim(0.0, 1.1)
    ax_right.set_yscale("log")
    ax_left.set_xscale("log")
    ax_left.set_xticks(utility_scales)
    ax_left.set_xticklabels([f"{value:g}" for value in utility_scales])
    ax_left.grid(True, which="both", linestyle="--", alpha=0.35)

    lines = [line_overlap, line_pearson, line_nsr]
    labels = [line.get_label() for line in lines]
    ax_left.legend(lines, labels, loc="lower left", frameon=True)

    fig.tight_layout()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(output_path, bbox_inches="tight")
    plt.close(fig)


def cleanup_old_evaluator_figures(result_dir: Path) -> None:
    old_filenames = [
        "noise_signal_ratio_curve.png",
        "overlap_at_1_curve.png",
        "overlap_at_3_curve.png",
        "overlap_at_5_curve.png",
        "overlap_at_10_curve.png",
        "pearson_correlation_curve.png",
        "mean_absolute_drift_curve.png",
        "max_absolute_drift_curve.png",
        "mrr_at_5_curve.png",
        "direction_cosine_curve.png",
        "retrieval_time_curve.png",
        "dp_noise_time_curve.png",
    ]
    for filename in old_filenames:
        path = result_dir / filename
        if path.is_file():
            path.unlink()


def evaluate(args: argparse.Namespace) -> None:
    os.makedirs(RESULT_DIR, exist_ok=True)
    cleanup_old_evaluator_figures(RESULT_DIR)
    context = prepare_evaluation_context(args)
    results = [
        evaluate_utility_scale(args, context, utility_scale)
        for utility_scale in UTILITY_SCALE_LIST
    ]
    output_path = RESULT_DIR / "privacy_utility_tradeoff.png"
    plot_privacy_utility_tradeoff(results, output_path)
    print("\nSaved evaluation figures:")
    print(f"  {output_path}")


def main() -> None:
    evaluate(parse_args())


if __name__ == "__main__":
    main()
