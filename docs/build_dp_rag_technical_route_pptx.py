"""Build a fully editable PowerPoint technical-route and principle diagram."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).with_name("DP_RAG_技术路线与原理图_可编辑.pptx")

# Modern academic palette: warm paper, navy ink, restrained teal/coral/gold.
BG = "F7F8FA"
WHITE = "FFFFFF"
INK = "17233B"
MUTED = "64748B"
LINE = "D9E1EA"
TEAL = "1B8A89"
TEAL_LIGHT = "E3F3F1"
BLUE = "3D6FB4"
BLUE_LIGHT = "E9F0FA"
CORAL = "D9685A"
CORAL_LIGHT = "FBEAE7"
GOLD = "D99A32"
GOLD_LIGHT = "FFF3D8"
GREEN = "4B9560"
GREEN_LIGHT = "E9F5EC"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_text(slide, text, x, y, w, h, size=16, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Microsoft YaHei", margin=0.05,
             valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(0.01)
    box.text_frame.margin_bottom = Inches(0.01)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def shape(slide, kind, x, y, w, h, fill=WHITE, line=LINE, radius=True, lw=1.2):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(line); s.line.width = Pt(lw)
    return s


def card(slide, x, y, w, h, fill=WHITE, line=LINE, accent=None):
    s = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line, lw=1.1)
    if accent:
        a = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.08, h, accent, accent, lw=0)
        a.adjustments[0] = 0.12
    return s


def pill(slide, text, x, y, w, fill, color, size=10):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.3, fill, fill, lw=0)
    return add_text(slide, text, x, y, w, 0.3, size, color, True, PP_ALIGN.CENTER)


def connector(slide, x1, y1, x2, y2, color=INK, width=1.8, dash=False, arrow=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = rgb(color); c.line.width = Pt(width)
    if dash: c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow: c.line.end_arrowhead = True
    return c


def dot(slide, x, y, d=0.12, fill=TEAL, line=None):
    return shape(slide, MSO_SHAPE.OVAL, x, y, d, d, fill, line or fill, lw=0.6)


def icon_document(slide, x, y, color=BLUE):
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.38, 0.5, WHITE, color, lw=1.3)
    shape(slide, MSO_SHAPE.FOLDED_CORNER, x + 0.08, y - 0.06, 0.38, 0.5, WHITE, color, lw=1.3)
    for i in range(3): connector(slide, x + 0.15, y + 0.1 + i*0.1, x + 0.37, y + 0.1 + i*0.1, color, 0.8, arrow=False)


def icon_shield(slide, x, y, color=CORAL):
    s = shape(slide, MSO_SHAPE.PENTAGON, x, y, 0.54, 0.58, CORAL_LIGHT, color, lw=1.5)
    s.rotation = 180
    add_text(slide, "✓", x, y + 0.02, 0.54, 0.45, 18, color, True, PP_ALIGN.CENTER)


def icon_vectors(slide, x, y):
    for i, color in enumerate((BLUE, TEAL, GOLD)):
        connector(slide, x, y + 0.46, x + 0.30 + i*0.11, y + 0.12 + i*0.08, color, 2.0, arrow=True)
    dot(slide, x - 0.05, y + 0.41, 0.11, INK)


def icon_graph(slide, x, y):
    pts = [(x, y+.26), (x+.32, y), (x+.38, y+.48), (x+.7, y+.18), (x+.75, y+.52)]
    edges = [(0,1),(0,2),(1,3),(2,3),(2,4),(3,4)]
    for a,b in edges: connector(slide, *pts[a], *pts[b], MUTED, 1.0, arrow=False)
    for i,(px,py) in enumerate(pts): dot(slide, px-.06, py-.06, .13, TEAL if i in (2,3) else BLUE_LIGHT, TEAL)


def icon_llm(slide, x, y):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.76, 0.56, BLUE_LIGHT, BLUE, lw=1.4)
    add_text(slide, "LLM", x, y, 0.76, 0.56, 14, BLUE, True, PP_ALIGN.CENTER)
    for yy in (y+.12, y+.42):
        dot(slide, x-.09, yy, .09, GOLD); dot(slide, x+.76, yy, .09, GOLD)


def section_label(slide, text, x, y, color):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.22, 0.32, color, color, lw=0)
    add_text(slide, text, x, y, 1.22, 0.32, 11, WHITE, True, PP_ALIGN.CENTER)


def add_header(slide, index, title, subtitle):
    add_text(slide, f"0{index}", 0.55, 0.34, 0.55, 0.4, 13, TEAL, True)
    add_text(slide, title, 1.05, 0.27, 8.6, 0.55, 25, INK, True)
    add_text(slide, subtitle, 9.55, 0.34, 3.15, 0.4, 10, MUTED, False, PP_ALIGN.RIGHT)
    connector(slide, 0.55, 0.88, 12.78, 0.88, LINE, 1.0, arrow=False)


def build_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
    add_header(slide, 1, "隐私增强 DP‑RAG：技术路线", "可编辑原理架构 · 离线保护 / 在线问答")

    # Client ribbon
    card(slide, .55, 1.08, 12.23, .66, WHITE, LINE)
    pill(slide, "客户端", .75, 1.25, .92, BLUE_LIGHT, BLUE)
    add_text(slide, "用户查询 q", 1.9, 1.2, 1.45, .35, 15, INK, True)
    connector(slide, 3.35, 1.40, 4.05, 1.40, TEAL, 2.2)
    add_text(slide, "语义检索", 4.15, 1.2, 1.25, .35, 13, TEAL, True, PP_ALIGN.CENTER)
    connector(slide, 5.45, 1.40, 6.15, 1.40, TEAL, 2.2)
    add_text(slide, "隐私知识上下文", 6.22, 1.2, 1.75, .35, 13, CORAL, True, PP_ALIGN.CENTER)
    connector(slide, 8.05, 1.40, 8.72, 1.40, TEAL, 2.2)
    add_text(slide, "本地 LLM", 8.82, 1.2, 1.3, .35, 13, BLUE, True, PP_ALIGN.CENTER)
    connector(slide, 10.15, 1.40, 10.8, 1.40, TEAL, 2.2)
    add_text(slide, "安全回答", 10.92, 1.2, 1.25, .35, 15, INK, True)

    # Offline lane
    section_label(slide, "离线 · 隐私建库", .55, 1.98, INK)
    # 1 docs
    card(slide, .55, 2.38, 1.7, 1.45, WHITE, LINE, BLUE)
    icon_document(slide, .78, 2.66)
    add_text(slide, "私有知识库", 1.35, 2.52, .72, .32, 13, INK, True)
    add_text(slide, "多格式加载\n重叠分块", .78, 3.22, 1.25, .45, 10, MUTED)
    # 2 scoring
    card(slide, 2.55, 2.18, 2.15, 1.85, WHITE, LINE, CORAL)
    icon_shield(slide, 2.8, 2.49)
    add_text(slide, "混合敏感度感知", 3.45, 2.4, 1.1, .35, 13, INK, True)
    pill(slide, "正则", 2.8, 3.04, .48, CORAL_LIGHT, CORAL, 9)
    pill(slide, "关键词", 3.34, 3.04, .62, CORAL_LIGHT, CORAL, 9)
    pill(slide, "语义", 4.02, 3.04, .48, CORAL_LIGHT, CORAL, 9)
    add_text(slide, "输出：敏感分数 sᵢ", 2.8, 3.48, 1.65, .3, 11, CORAL, True)
    # 3 representation
    card(slide, 5.0, 2.38, 1.85, 1.45, WHITE, LINE, TEAL)
    icon_vectors(slide, 5.25, 2.58)
    add_text(slide, "统一表征", 5.88, 2.48, .78, .32, 13, INK, True)
    add_text(slide, "L2 归一化\nJL 投影 d → m", 5.25, 3.17, 1.4, .45, 10, MUTED)
    # central principle engine
    card(slide, 7.15, 2.06, 3.05, 2.1, CORAL_LIGHT, CORAL, CORAL)
    pill(slide, "核心机制", 7.43, 2.33, .82, CORAL, WHITE)
    add_text(slide, "敏感度驱动的\n动态解析高斯机制", 8.38, 2.25, 1.58, .62, 16, INK, True)
    add_text(slide, "高敏感", 7.47, 3.12, .55, .26, 10, CORAL, True)
    connector(slide, 8.12, 3.25, 9.18, 3.25, CORAL, 2.0)
    add_text(slide, "低 εᵢ  ·  大 Δᵢ", 9.24, 3.10, .75, .32, 10, CORAL, True, PP_ALIGN.RIGHT)
    add_text(slide, "zᵢ ~ 𝒩(0, (σᵢ·α/√m)²I)", 7.44, 3.57, 2.45, .3, 12, INK, True, PP_ALIGN.CENTER)
    # 5 index
    card(slide, 10.5, 2.38, 2.28, 1.45, WHITE, LINE, GREEN)
    icon_graph(slide, 10.8, 2.55)
    add_text(slide, "隐私向量索引", 11.62, 2.48, .95, .32, 13, INK, True)
    add_text(slide, "最终归一化\nHNSW 建图", 10.82, 3.18, 1.45, .45, 10, MUTED)
    # offline connectors
    for x1,x2 in ((2.25,2.55),(4.70,5.0),(6.85,7.15),(10.2,10.5)):
        connector(slide, x1, 3.08, x2, 3.08, TEAL if x1 < 6.85 else CORAL, 2.1)

    # Online lane and components
    section_label(slide, "在线 · 检索生成", .55, 4.45, TEAL)
    card(slide, .55, 4.86, 2.15, 1.28, TEAL_LIGHT, TEAL)
    add_text(slide, "① 查询同构映射", .78, 5.05, 1.6, .3, 13, INK, True)
    add_text(slide, "q → Embedding → JL", .78, 5.48, 1.55, .28, 11, TEAL, True)
    add_text(slide, "与文档共享同一投影矩阵", .78, 5.78, 1.7, .22, 9, MUTED)
    card(slide, 3.0, 4.86, 2.15, 1.28, GREEN_LIGHT, GREEN)
    icon_graph(slide, 3.26, 5.18)
    add_text(slide, "② HNSW Top‑K", 4.03, 5.03, .92, .3, 13, INK, True)
    add_text(slide, "近似近邻检索", 4.02, 5.50, .9, .26, 10, GREEN, True)
    card(slide, 5.45, 4.86, 2.25, 1.28, GOLD_LIGHT, GOLD)
    add_text(slide, "③ 上下文构建", 5.73, 5.05, 1.45, .3, 13, INK, True)
    pill(slide, "Top‑K 片段", 5.73, 5.48, .86, WHITE, GOLD, 9)
    pill(slide, "Identity", 6.66, 5.48, .72, WHITE, GOLD, 9)
    add_text(slide, "只暴露必要的检索结果", 5.73, 5.82, 1.55, .2, 9, MUTED)
    card(slide, 8.0, 4.86, 2.15, 1.28, BLUE_LIGHT, BLUE)
    icon_llm(slide, 8.27, 5.16)
    add_text(slide, "④ 本地生成", 9.15, 5.07, .75, .3, 13, INK, True)
    add_text(slide, "知识增强回答", 9.12, 5.52, .8, .25, 10, BLUE, True)
    card(slide, 10.45, 4.86, 2.33, 1.28, WHITE, LINE, INK)
    add_text(slide, "⑤ 输出与审计", 10.72, 5.05, 1.42, .3, 13, INK, True)
    add_text(slide, "回答 + 检索来源\n隐私 / 效用指标", 10.72, 5.48, 1.55, .48, 10, MUTED)
    for x1,x2 in ((2.7,3.0),(5.15,5.45),(7.7,8.0),(10.15,10.45)):
        connector(slide, x1, 5.5, x2, 5.5, TEAL, 2.1)
    # cross-lane logic arrows
    connector(slide, 11.65, 3.83, 4.1, 4.83, GREEN, 1.3, dash=True)
    add_text(slide, "检索索引", 7.43, 4.2, .8, .25, 9, GREEN, True, PP_ALIGN.CENTER)
    connector(slide, 6.0, 3.83, 1.63, 4.83, TEAL, 1.3, dash=True)
    add_text(slide, "共享投影 R", 3.38, 4.18, .95, .25, 9, TEAL, True, PP_ALIGN.CENTER)

    # Footer principles
    connector(slide, .55, 6.48, 12.78, 6.48, LINE, 1, arrow=False)
    pill(slide, "隐私", .62, 6.72, .55, CORAL_LIGHT, CORAL)
    add_text(slide, "分块级动态 (εᵢ, δ)-DP", 1.25, 6.69, 2.05, .3, 10, MUTED)
    pill(slide, "效用", 3.6, 6.72, .55, TEAL_LIGHT, TEAL)
    add_text(slide, "JL 保距 + 维度校正噪声", 4.23, 6.69, 2.15, .3, 10, MUTED)
    pill(slide, "效率", 6.75, 6.72, .55, GREEN_LIGHT, GREEN)
    add_text(slide, "HNSW 亚线性近似检索", 7.38, 6.69, 1.95, .3, 10, MUTED)
    pill(slide, "边界", 9.72, 6.72, .55, BLUE_LIGHT, BLUE)
    add_text(slide, "原文与 LLM 均保持本地", 10.35, 6.69, 2.1, .3, 10, MUTED)
    return slide


def build_principle(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(BG)
    add_header(slide, 2, "核心原理：敏感度如何驱动噪声", "从文本风险到分块级隐私预算")

    # left scoring funnel
    card(slide, .55, 1.2, 3.05, 5.65, WHITE, LINE, CORAL)
    add_text(slide, "A. 混合敏感度评估", .85, 1.48, 2.4, .4, 18, INK, True)
    icon_shield(slide, 2.93, 1.42)
    rows = [("直接标识符", "身份证 / 手机 / 银行卡", CORAL), ("隐私关键词", "个人、财务、企业机密", GOLD), ("上下文语义", "启发式 + 零样本分类", BLUE)]
    y = 2.15
    for name,desc,col in rows:
        card(slide, .85, y, 2.45, .78, BG, LINE)
        dot(slide, 1.03, y+.31, .14, col)
        add_text(slide, name, 1.25, y+.12, 1.5, .25, 12, INK, True)
        add_text(slide, desc, 1.25, y+.41, 1.75, .22, 9, MUTED)
        y += .95
    connector(slide, 2.07, 4.92, 2.07, 5.35, CORAL, 2.2)
    shape(slide, MSO_SHAPE.OVAL, 1.33, 5.35, 1.48, .88, CORAL_LIGHT, CORAL, lw=1.5)
    add_text(slide, "敏感分数\nsᵢ ∈ [0,1]", 1.43, 5.46, 1.28, .62, 15, CORAL, True, PP_ALIGN.CENTER)
    add_text(slide, "分块级，而非整库统一", 1.05, 6.42, 2.05, .24, 10, MUTED, False, PP_ALIGN.CENTER)

    # center mapping
    card(slide, 3.95, 1.2, 4.18, 5.65, CORAL_LIGHT, CORAL, CORAL)
    add_text(slide, "B. 动态隐私预算与局部敏感度", 4.28, 1.48, 3.5, .4, 18, INK, True)
    # axes diagram, editable lines/dots
    add_text(slide, "εᵢ", 4.42, 2.12, .35, .3, 11, CORAL, True)
    connector(slide, 4.72, 4.12, 4.72, 2.3, MUTED, 1.2, arrow=True)
    connector(slide, 4.72, 4.12, 7.55, 4.12, MUTED, 1.2, arrow=True)
    add_text(slide, "敏感度 sᵢ", 6.72, 4.22, .85, .25, 10, MUTED)
    pts = [(4.87,2.52),(5.35,2.78),(5.85,3.12),(6.35,3.49),(6.92,3.82)]
    for a,b in zip(pts,pts[1:]): connector(slide,*a,*b,CORAL,2.4,arrow=False)
    for px,py in pts: dot(slide,px-.06,py-.06,.12,CORAL)
    add_text(slide, "低敏感：ε 大，保护较弱", 4.72, 4.65, 2.7, .28, 11, TEAL, True)
    add_text(slide, "高敏感：ε 小，保护更强", 4.72, 5.02, 2.7, .28, 11, CORAL, True)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 4.48, 5.55, 3.18, .78, WHITE, LINE, lw=1)
    add_text(slide, "εᵢ = 1.25 + 8.75(1−sᵢ)¹·⁵\nΔᵢ = 0.25 + 0.25sᵢ", 4.65, 5.64, 2.82, .58, 13, INK, True, PP_ALIGN.CENTER)

    # right noise and utility
    card(slide, 8.48, 1.2, 4.3, 5.65, WHITE, LINE, TEAL)
    add_text(slide, "C. 解析高斯校准与效用控制", 8.8, 1.48, 3.55, .4, 18, INK, True)
    add_text(slide, "解析求解", 8.86, 2.12, .9, .25, 11, CORAL, True)
    add_text(slide, "g(u*) = δ", 9.85, 2.05, 1.35, .4, 17, INK, True, PP_ALIGN.CENTER)
    connector(slide, 11.22, 2.25, 11.72, 2.25, CORAL, 2.0)
    add_text(slide, "σᵢ = u*Δᵢ", 11.72, 2.05, .82, .4, 13, CORAL, True, PP_ALIGN.RIGHT)
    # bell curves made of dots/segments
    connector(slide, 8.9, 4.18, 12.35, 4.18, MUTED, 1.0, arrow=True)
    connector(slide, 10.62, 4.25, 10.62, 2.72, MUTED, 1.0, arrow=True)
    narrow=[(9.55,4.12),(9.9,3.93),(10.18,3.42),(10.42,2.94),(10.62,2.79),(10.82,2.94),(11.06,3.42),(11.34,3.93),(11.69,4.12)]
    wide=[(9.15,4.1),(9.52,3.9),(9.87,3.6),(10.23,3.3),(10.62,3.18),(11.01,3.3),(11.37,3.6),(11.72,3.9),(12.09,4.1)]
    for seq,col in ((narrow,TEAL),(wide,CORAL)):
        for a,b in zip(seq,seq[1:]): connector(slide,*a,*b,col,2.1,arrow=False)
    pill(slide, "低敏感 / 小噪声", 8.92, 4.52, 1.25, TEAL_LIGHT, TEAL, 9)
    pill(slide, "高敏感 / 大噪声", 10.35, 4.52, 1.32, CORAL_LIGHT, CORAL, 9)
    add_text(slide, "维度校正避免高维噪声能量爆炸", 8.92, 5.12, 3.35, .28, 11, INK, True, PP_ALIGN.CENTER)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 9.05, 5.55, 3.08, .78, TEAL_LIGHT, TEAL, lw=1)
    add_text(slide, "σᵢ,dim = σᵢ·α / √m\nzᵢ ~ 𝒩(0, σ²ᵢ,dim I)", 9.25, 5.64, 2.68, .58, 13, INK, True, PP_ALIGN.CENTER)

    connector(slide, 3.6, 4.0, 3.95, 4.0, CORAL, 2.3)
    connector(slide, 8.13, 4.0, 8.48, 4.0, CORAL, 2.3)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_overview(prs)
    build_principle(prs)
    prs.core_properties.title = "隐私增强 DP-RAG 技术路线与原理图"
    prs.core_properties.subject = "全部由 PowerPoint 原生可编辑形状构成"
    prs.core_properties.author = "Codex"
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
