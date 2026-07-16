"""Editable second slide: sensitivity-driven dynamic DP mechanism."""

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_dp_rag_technical_route_pptx import add_text, connector, dot, icon_graph, icon_shield, rgb, shape

OUT=Path(__file__).with_name("DP_RAG_第二页_动态差分隐私核心原理_可编辑.pptx")
INK="17202B"; MUTED="5D6875"; BG="FFFFFF"; LINE="D2DAE3"
BLUE="2468B2"; BLUE_L="F2F7FD"; GREEN="267A3F"; GREEN_L="F1F8F2"
PURPLE="65379F"; PURPLE_L="F6F1FB"; ORANGE="C75B1C"; ORANGE_L="FFF4EC"
RED="B33A3A"; RED_L="FFF0F0"; GOLD="D59B31"; GOLD_L="FFF7E4"

def card(s,x,y,w,h,fill=BG,line=LINE,lw=.9): return shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,fill,line,lw=lw)
def arrow(s,x1,y1,x2,y2,color=GREEN,width=1.3): return connector(s,x1,y1,x2,y2,color,width,arrow=True)
def badge(s,n,x,y,color=GREEN):
    shape(s,MSO_SHAPE.OVAL,x,y,.25,.25,color,color,lw=0)
    add_text(s,str(n),x,y,.25,.25,8,"FFFFFF",True,PP_ALIGN.CENTER)
def title_card(s,title,x,y,w,h,color,fill):
    card(s,x,y,w,h,fill,color,1.0); add_text(s,title,x+.12,y+.05,w-.24,.28,12,color,True,PP_ALIGN.CENTER)
def chunk(s,x,y,color,label):
    shape(s,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,.58,.34,"FFFFFF",color,lw=.9)
    add_text(s,label,x,y,.58,.34,8.5,color,True,PP_ALIGN.CENTER)
def bell(s,x,y,width,color):
    pts=[(x,y+.48),(x+width*.17,y+.40),(x+width*.34,y+.22),(x+width*.50,y),(x+width*.66,y+.22),(x+width*.83,y+.40),(x+width,y+.48)]
    for a,b in zip(pts,pts[1:]): connector(s,*a,*b,color,1.5,arrow=False)

def build():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(BG)
    add_text(s,"敏感度驱动的动态差分隐私保护机制",1.55,.10,10.20,.38,20,INK,True,PP_ALIGN.CENTER)
    add_text(s,"根据文档分块的隐私风险动态调整隐私预算、局部敏感度和噪声强度，实现差异化语义保护与检索效用控制",1.15,.49,11.0,.25,10.5,MUTED,False,PP_ALIGN.CENTER)
    connector(s,.20,.82,13.13,.82,BLUE,1.1,arrow=False)

    # Main left 80% and right 20% frames.
    card(s,.20,.98,10.25,5.43,BLUE_L,BLUE,1.05)
    card(s,10.62,.98,2.51,5.43,PURPLE_L,PURPLE,1.05)
    add_text(s,"动态隐私保护原理",3.82,1.05,3.0,.28,14,BLUE,True,PP_ALIGN.CENTER)
    add_text(s,"方法特性",11.12,1.05,1.52,.28,14,PURPLE,True,PP_ALIGN.CENTER)

    # Upper mechanism strip: three large panels.
    title_card(s,"① 分块级敏感度评估",.42,1.45,2.80,2.15,BLUE, "FFFFFF")
    chunk(s,.65,1.98,BLUE,"分块 1"); chunk(s,1.34,1.98,GOLD,"分块 2"); chunk(s,2.03,1.98,RED,"分块 3")
    add_text(s,"…",2.67,2.02,.25,.24,12,INK,True,PP_ALIGN.CENTER)
    arrow(s,1.82,2.38,1.82,2.62,BLUE,1.0)
    card(s,.67,2.65,2.28,.40,"F8FAFC",LINE)
    add_text(s,"规则匹配 + 关键词 + 语义判断",.69,2.65,2.24,.40,9,INK,True,PP_ALIGN.CENTER)
    # sensitivity ruler
    connector(s,.82,3.30,2.78,3.30,"68A76B",2.2,arrow=False)
    connector(s,1.78,3.30,2.78,3.30,RED,2.2,arrow=False)
    dot(s,2.15,3.23,.14,INK,INK)
    add_text(s,"低敏感",.68,3.34,.58,.18,7.5,GREEN,True,PP_ALIGN.CENTER)
    add_text(s,"高敏感",2.38,3.34,.58,.18,7.5,RED,True,PP_ALIGN.CENTER)
    add_text(s,"原始分数 rᵢ → 归一化敏感度 sᵢ",.77,3.07,2.10,.20,8.2,MUTED,False,PP_ALIGN.CENTER)

    title_card(s,"② 动态隐私参数映射",3.43,1.45,3.18,2.15,GREEN,"FFFFFF")
    # editable axes and two curves
    connector(s,3.78,3.12,3.78,1.93,MUTED,.8,arrow=True); connector(s,3.78,3.12,6.18,3.12,MUTED,.8,arrow=True)
    add_text(s,"sᵢ",6.00,3.14,.23,.18,8,MUTED,True,PP_ALIGN.CENTER)
    eps=[(3.93,2.08),(4.35,2.18),(4.78,2.38),(5.22,2.66),(5.73,2.94)]
    delta=[(3.93,2.91),(4.35,2.78),(4.78,2.61),(5.22,2.38),(5.73,2.12)]
    for seq,col in ((eps,RED),(delta,GREEN)):
        for a,b in zip(seq,seq[1:]): connector(s,*a,*b,col,1.7,arrow=False)
    add_text(s,"εᵢ ↓",5.68,2.93,.42,.19,8,RED,True,PP_ALIGN.CENTER)
    add_text(s,"Δᵢ ↑",5.68,2.00,.42,.19,8,GREEN,True,PP_ALIGN.CENTER)
    add_text(s,"εᵢ = 1.25 + 8.75(1−sᵢ)¹·⁵",3.92,3.27,2.30,.18,8.5,INK,True,PP_ALIGN.CENTER)
    add_text(s,"Δᵢ = 0.25 + 0.25sᵢ",4.18,3.45,1.78,.17,8.5,INK,True,PP_ALIGN.CENTER)

    title_card(s,"③ 差异化高斯噪声分配",6.82,1.45,3.38,2.15,GREEN,"FFFFFF")
    bell(s,7.04,2.38,.72,BLUE); bell(s,7.89,2.20,1.00,GOLD); bell(s,9.00,2.02,1.00,RED)
    add_text(s,"低敏感",6.94,2.95,.90,.20,8,BLUE,True,PP_ALIGN.CENTER)
    add_text(s,"中敏感",7.95,2.95,.90,.20,8,GOLD,True,PP_ALIGN.CENTER)
    add_text(s,"高敏感",9.04,2.95,.90,.20,8,RED,True,PP_ALIGN.CENTER)
    add_text(s,"小噪声",6.94,3.16,.90,.18,7.5,MUTED,False,PP_ALIGN.CENTER)
    add_text(s,"适中噪声",7.95,3.16,.90,.18,7.5,MUTED,False,PP_ALIGN.CENTER)
    add_text(s,"强噪声",9.04,3.16,.90,.18,7.5,MUTED,False,PP_ALIGN.CENTER)
    add_text(s,"σᵢ = u*Δᵢ　　zᵢ ~ N(0, σ²ᵢ,dim I)",7.09,3.39,2.85,.18,8.5,INK,True,PP_ALIGN.CENTER)
    arrow(s,3.22,2.48,3.43,2.48,GREEN,1.4); arrow(s,6.61,2.48,6.82,2.48,GREEN,1.4)

    # Lower processing chain.
    add_text(s,"向量保护与检索效用控制",3.70,3.79,3.25,.27,13,BLUE,True,PP_ALIGN.CENTER)
    lower=[
        ("原始语义向量","vᵢ",BLUE),
        ("L2 范数裁剪","限制向量范数 C",GREEN),
        ("维度能量修正","σᵢ,dim = σᵢ·α/√d",ORANGE),
        ("高斯噪声注入","vᵢᶜ + zᵢ",RED),
        ("最终归一化","输出单位向量",GREEN),
        ("HNSW 索引","余弦 Top‑K",PURPLE),
    ]
    x=.48
    for i,(name,note,col) in enumerate(lower):
        w=1.42 if i<5 else 1.55
        card(s,x,4.24,w,.96,"FFFFFF",col)
        badge(s,i+1,x+.08,4.34,col)
        add_text(s,name,x+.32,4.31,w-.38,.22,9.5,col,True,PP_ALIGN.CENTER)
        add_text(s,note,x+.10,4.66,w-.20,.28,8,INK,False,PP_ALIGN.CENTER)
        if i<len(lower)-1: arrow(s,x+w,4.72,x+w+.20,4.72,BLUE,1.1)
        x+=w+.20
    # Explanatory callouts.
    card(s,.48,5.38,3.00,.67,"FFFFFF",LINE)
    add_text(s,"裁剪作用",.62,5.44,.72,.20,9,GREEN,True)
    add_text(s,"建立统一敏感度边界，为高斯机制提供稳定输入。",1.25,5.43,2.08,.42,8.2,MUTED,False)
    card(s,3.64,5.38,3.18,.67,"FFFFFF",LINE)
    add_text(s,"维度修正",3.78,5.44,.82,.20,9,ORANGE,True)
    add_text(s,"抑制高维噪声能量累积，减少相似度结构破坏。",4.53,5.43,2.14,.42,8.2,MUTED,False)
    card(s,6.98,5.38,3.20,.67,"FFFFFF",LINE)
    add_text(s,"检索适配",7.12,5.44,.72,.20,9,PURPLE,True)
    add_text(s,"归一化后构建 HNSW，在隐私约束下保持检索效率。",7.78,5.43,2.24,.42,8.2,MUTED,False)

    # Right feature cards.
    features=[
        ("差异化保护","不同敏感度分块采用不同隐私预算和噪声强度，避免统一加噪造成保护不足或过度扰动。",GREEN,GREEN_L),
        ("效用控制","结合 L2 裁剪、维度能量修正和最终归一化，降低噪声对向量方向的破坏。",BLUE,BLUE_L),
        ("高效检索","使用 HNSW 构建隐私向量索引，在数据规模增长时保持较低查询开销。",PURPLE,PURPLE_L),
    ]
    y=1.47
    for i,(name,desc,col,fill) in enumerate(features,1):
        card(s,10.84,y,2.07,1.18,"FFFFFF",col)
        badge(s,i,10.96,y+.12,col)
        add_text(s,name,11.30,y+.10,1.35,.23,11,col,True,PP_ALIGN.CENTER)
        add_text(s,desc,11.02,y+.39,1.73,.65,8.4,INK,False,PP_ALIGN.LEFT)
        y+=1.34
    card(s,10.84,5.54,2.07,.63,ORANGE_L,ORANGE)
    add_text(s,"核心贡献",11.32,5.59,1.15,.20,10,ORANGE,True,PP_ALIGN.CENTER)
    add_text(s,"将保护强度从“全局固定”转变为“分块自适应”",11.02,5.82,1.72,.28,8.4,INK,True,PP_ALIGN.CENTER)

    # Bottom comparison band, same thin-card language as page one.
    card(s,.20,6.56,12.93,.73,RED_L,"E2A1A1",.8)
    add_text(s,"统一加噪",.44,6.66,.86,.20,9.5,RED,True,PP_ALIGN.CENTER)
    add_text(s,"固定预算 · 固定噪声 · 高敏感内容可能保护不足 · 低敏感内容容易过度扰动",1.34,6.65,4.62,.24,8.5,INK,False,PP_ALIGN.CENTER)
    shape(s,MSO_SHAPE.CHEVRON,6.10,6.70,.38,.24,RED,RED,lw=0)
    add_text(s,"动态加噪",6.66,6.66,.86,.20,9.5,GREEN,True,PP_ALIGN.CENTER)
    add_text(s,"动态预算 · 差异化噪声 · 强化敏感内容保护 · 更好保留低敏感内容的检索效用",7.54,6.65,5.28,.24,8.5,INK,False,PP_ALIGN.CENTER)

    prs.core_properties.title="敏感度驱动的动态差分隐私保护机制"
    prs.core_properties.subject="与隐私保护 RAG 系统架构页配套的可编辑第二页"
    prs.core_properties.author="Codex"
    prs.save(OUT); print(OUT)

if __name__=="__main__": build()
