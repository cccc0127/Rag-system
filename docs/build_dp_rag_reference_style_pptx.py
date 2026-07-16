"""Create an editable Chinese DP-RAG principle diagram in the reference-paper style."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from build_dp_rag_technical_route_pptx import (
    add_text, connector, dot, icon_document, icon_graph, icon_llm, icon_shield,
    rgb, shape,
)


OUT = Path(__file__).with_name("DP_RAG_中文版原理图_参考风格_可编辑.pptx")
INK = "111111"
BLUE_BG = "DDECF8"
BLUE = "3976A8"
GREEN_BG = "E7F2DD"
GREEN = "56863C"
PURPLE_BG = "EDE1F5"
PURPLE = "764893"
GRAY = "E4E4E4"
PINK = "D95B9D"
CYAN = "5DB7C4"
ORANGE = "E9A343"
NAVY = "305A85"
WHITE = "FFFFFF"


def box(slide, x, y, w, h, fill="FFFFFF", line=INK, lw=2.1):
    return shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, lw=lw)


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=1.7):
    return connector(slide, x1, y1, x2, y2, color, width, arrow=True)


def tile_matrix(slide, x, y, rows=2, cols=5, cell=.15, palette=None, border="9AA5AE"):
    palette = palette or [PINK, "D9DEE2", PINK, "D9DEE2", CYAN, "A9D18E"]
    for r in range(rows):
        for c in range(cols):
            color = palette[(r * cols + c) % len(palette)]
            shape(slide, MSO_SHAPE.RECTANGLE, x+c*cell, y+r*cell, cell, cell, color, border, lw=.45)


def vector_bar(slide, x, y, count=12, color=ORANGE):
    for i in range(count):
        shape(slide, MSO_SHAPE.RECTANGLE, x+i*.055, y, .047, .20, WHITE, color, lw=.65)
    add_text(slide, "…", x+count*.055+.01, y-.04, .22, .22, 10, INK, True, PP_ALIGN.CENTER)


def number_caption(slide, num, text, x, y, w, color=INK, size=10.5):
    add_text(slide, f"{num}  {text}", x, y, w, .35, size, color, False, PP_ALIGN.CENTER)


def simple_embedding_model(slide, x, y):
    # Editable node-network icon.
    pts = [(x,y+.22),(x+.2,y),(x+.22,y+.45),(x+.46,y+.12),(x+.5,y+.42),(x+.72,y+.25)]
    for a,b in ((0,1),(0,2),(1,3),(2,3),(2,4),(3,5),(4,5)):
        connector(slide,*pts[a],*pts[b],NAVY,1,arrow=False)
    for i,(px,py) in enumerate(pts): dot(slide,px-.055,py-.055,.12,PINK if i in (1,4) else CYAN,NAVY)


def noise_bell(slide, x, y, wide=False, color=GREEN):
    if wide:
        pts=[(x,y+.42),(x+.18,y+.34),(x+.37,y+.2),(x+.58,y+.08),(x+.78,y+.2),(x+.97,y+.34),(x+1.15,y+.42)]
    else:
        pts=[(x,y+.42),(x+.25,y+.38),(x+.45,y+.19),(x+.57,y),(x+.69,y+.19),(x+.9,y+.38),(x+1.15,y+.42)]
    for a,b in zip(pts,pts[1:]): connector(slide,*a,*b,color,1.5,arrow=False)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(WHITE)

    # Phase I — reference-style narrow initialization band.
    add_text(slide, "阶段 I：离线初始化", 4.85, .05, 3.65, .38, 20, INK, True, PP_ALIGN.CENTER)
    box(slide, .02, .43, 3.83, .48, WHITE, INK, 1.8)
    add_text(slide, "私有知识库", .02, .47, 3.83, .36, 15, INK, True, PP_ALIGN.CENTER)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.85, .36, 5.63, .62, GRAY, "999999", lw=1.1)
    add_text(slide, "加载嵌入模型并生成共享 JL 投影矩阵 R", 3.98, .45, 5.36, .4, 15, INK, True, PP_ALIGN.CENTER)
    box(slide, 9.48, .43, 3.83, .48, WHITE, INK, 1.8)
    add_text(slide, "本地 RAG 服务端", 9.48, .47, 3.83, .36, 15, INK, True, PP_ALIGN.CENTER)

    add_text(slide, "阶段 II：隐私知识库构建", 4.65, .98, 4.0, .38, 20, INK, True, PP_ALIGN.CENTER)

    # Main pastel regions.
    box(slide, .02, 1.38, 5.25, 4.18, BLUE_BG, BLUE_BG, 0)
    box(slide, 5.27, 1.38, 8.04, 4.18, GREEN_BG, GREEN_BG, 0)
    box(slide, .02, 5.56, 13.29, 1.91, PURPLE_BG, PURPLE_BG, 0)
    add_text(slide, "1. 文档分块与隐私敏感度感知", .7, 1.46, 4.1, .38, 18, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, "2. 敏感度驱动的动态差分隐私表征", 6.9, 1.46, 4.65, .38, 18, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "3. 在线隐私检索与本地生成", 4.43, 5.61, 4.6, .38, 18, PURPLE, True, PP_ALIGN.CENTER)

    # Thick black subsystem outlines.
    box(slide, .32, 1.90, 4.63, 3.40, "F4FAFE", INK, 2.2)
    box(slide, 5.56, 1.90, 4.47, 3.40, "F7FBF2", INK, 2.2)
    box(slide, 10.22, 1.90, 2.78, 3.40, "F7FBF2", INK, 2.2)

    # Region 1: documents -> chunks -> scoring -> sensitivity vector.
    add_text(slide, "文档侧", 2.08, 1.95, 1.12, .32, 16, INK, True, PP_ALIGN.CENTER)
    icon_document(slide, .65, 2.38, NAVY)
    add_text(slide, "多格式文档", .47, 2.91, 1.12, .25, 10, INK, True, PP_ALIGN.CENTER)
    arrow(slide, 1.48, 2.62, 2.02, 2.62, BLUE)
    tile_matrix(slide, 2.08, 2.37, 3, 5, .17, [PINK,"D8E1E8",CYAN,"D8E1E8"])
    number_caption(slide,"1.1","重叠分块",1.93,2.94,1.23,NAVY)
    arrow(slide, 2.93, 2.70, 3.55, 2.70, BLUE)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.58, 2.31, .86, .78, "CCD9EA", "8099B6", lw=1.2)
    icon_shield(slide, 3.73, 2.39, NAVY)
    add_text(slide, "混合评估", 3.61, 3.14, .82, .25, 10, INK, True, PP_ALIGN.CENTER)
    # three scorer inputs
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,.62,3.58,1.03,.38,"FFFFFF","AAB9C4",lw=.8)
    add_text(slide,"正则匹配",.62,3.58,1.03,.38,10,INK,True,PP_ALIGN.CENTER)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,1.78,3.58,1.03,.38,"FFFFFF","AAB9C4",lw=.8)
    add_text(slide,"关键词规则",1.78,3.58,1.03,.38,10,INK,True,PP_ALIGN.CENTER)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,2.94,3.58,1.38,.38,"FFFFFF","AAB9C4",lw=.8)
    add_text(slide,"零样本语义分类",2.94,3.58,1.38,.38,9.5,INK,True,PP_ALIGN.CENTER)
    for xx in (1.13,2.29,3.63): arrow(slide,xx,3.57,3.91,3.10,BLUE,1.1)
    number_caption(slide,"1.2","融合规则分数与语义分数",1.05,4.12,3.25,NAVY)
    arrow(slide, 3.91, 3.40, 3.91, 4.35, BLUE)
    vector_bar(slide, 2.92, 4.48, 15, PINK)
    add_text(slide,"s₁  s₂  …  sₙ",2.96,4.76,1.35,.22,10,PINK,True,PP_ALIGN.CENTER)
    number_caption(slide,"1.3","输出分块级敏感度",1.02,4.87,1.85,NAVY)

    # Region 2A: shared representation and DP engine.
    add_text(slide, "隐私表征引擎", 6.92, 1.95, 1.78, .32, 16, INK, True, PP_ALIGN.CENTER)
    tile_matrix(slide, 5.85, 2.40, 3, 5, .16, [PINK,"D9DEE2",CYAN,"D9DEE2"])
    add_text(slide,"原始嵌入 xᵢ",5.72,2.94,1.18,.24,10,INK,True,PP_ALIGN.CENTER)
    arrow(slide,6.70,2.66,7.25,2.66,BLUE)
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,7.30,2.31,1.12,.72,"D8E4F2","819CB9",lw=1.1)
    add_text(slide,"L2 归一化\nJL 投影",7.30,2.35,1.12,.58,11,INK,True,PP_ALIGN.CENTER)
    arrow(slide,8.42,2.66,9.06,2.66,BLUE)
    tile_matrix(slide,9.12,2.45,2,4,.16,[CYAN,"D9DEE2",PINK,"D9DEE2"])
    add_text(slide,"低维向量 yᵢ",8.91,2.93,1.02,.24,10,INK,True,PP_ALIGN.CENTER)
    number_caption(slide,"2.1","共享表征：yᵢ = norm(xᵢR)",6.15,3.23,3.35,GREEN)
    # DP mechanism center lower
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,5.83,3.69,1.47,.80,"F5D7DD","C87587",lw=1.1)
    add_text(slide,"动态预算\nεᵢ ↓  Δᵢ ↑",5.83,3.75,1.47,.62,11,INK,True,PP_ALIGN.CENTER)
    arrow(slide,7.30,4.08,7.78,4.08,GREEN)
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,7.82,3.62,1.72,.94,"E6E6E6","999999",lw=1.1)
    add_text(slide,"解析高斯校准\nσᵢ = u*Δᵢ",7.82,3.70,1.72,.62,11,INK,True,PP_ALIGN.CENTER)
    arrow(slide,9.54,4.08,9.91,4.08,GREEN)
    noise_bell(slide,8.08,4.70,False,BLUE)
    noise_bell(slide,8.28,4.70,True,PINK)
    add_text(slide,"zᵢ ~ 𝒩(0,(σᵢα/√m)²I)",6.02,4.72,1.95,.28,10,INK,True,PP_ALIGN.CENTER)
    number_caption(slide,"2.2","高敏感 → 小 εᵢ → 强噪声",6.2,5.00,3.35,GREEN)
    # score cross-region input
    arrow(slide,4.46,4.62,5.78,4.15,PINK,1.8)

    # Region 2B: index creation.
    add_text(slide, "隐私索引", 11.02, 1.95, 1.18, .32, 16, INK, True, PP_ALIGN.CENTER)
    tile_matrix(slide,10.55,2.42,3,5,.16,[PINK,"D9DEE2",CYAN,"B9D692",PINK])
    add_text(slide,"扰动向量 ỹᵢ",10.47,2.98,1.03,.24,10,INK,True,PP_ALIGN.CENTER)
    arrow(slide,11.43,2.75,11.43,3.35,GREEN)
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,10.83,3.36,1.22,.52,"FFFFFF","9CB18E",lw=1)
    add_text(slide,"最终归一化",10.83,3.40,1.22,.38,10,INK,True,PP_ALIGN.CENTER)
    arrow(slide,11.43,3.88,11.43,4.18,GREEN)
    icon_graph(slide,10.88,4.26)
    add_text(slide,"HNSW 图索引",10.61,4.88,1.68,.25,11,GREEN,True,PP_ALIGN.CENTER)
    number_caption(slide,"2.3","构建隐私向量库",10.47,5.02,1.93,GREEN)
    arrow(slide,10.03,4.05,10.50,2.72,GREEN,1.8)

    # Phase III bottom lane: compact but illustrated.
    box(slide,.34,6.04,12.64,1.18,"F8F1FC",INK,2.0)
    shape(slide,MSO_SHAPE.OVAL,.62,6.31,.48,.48,"FFFFFF",PURPLE,lw=1.2)
    add_text(slide,"问",.62,6.31,.48,.48,13,PURPLE,True,PP_ALIGN.CENTER)
    add_text(slide,"用户查询 q",1.18,6.35,.85,.28,11,INK,True,PP_ALIGN.CENTER)
    arrow(slide,2.08,6.56,2.63,6.56,PURPLE)
    simple_embedding_model(slide,2.72,6.30)
    add_text(slide,"嵌入 + 同一 JL",2.57,6.84,1.15,.22,9.5,INK,True,PP_ALIGN.CENTER)
    arrow(slide,3.82,6.56,4.40,6.56,PURPLE)
    icon_graph(slide,4.56,6.30)
    add_text(slide,"HNSW Top‑K",4.42,6.84,1.23,.22,9.5,INK,True,PP_ALIGN.CENTER)
    arrow(slide,5.73,6.56,6.30,6.56,PURPLE)
    # context sheets
    for i,col in enumerate(("D7E5F1","F1D8E7","DDEBCF")):
        shape(slide,MSO_SHAPE.RECTANGLE,6.41+i*.08,6.27-i*.04,.65,.50,col,"8B8B8B",lw=.7)
    add_text(slide,"检索片段 + 提示词",6.15,6.84,1.55,.22,9.5,INK,True,PP_ALIGN.CENTER)
    arrow(slide,7.45,6.56,8.05,6.56,PURPLE)
    icon_llm(slide,8.18,6.27)
    add_text(slide,"本地 LLM 生成",8.02,6.84,1.15,.22,9.5,INK,True,PP_ALIGN.CENTER)
    arrow(slide,9.13,6.56,9.75,6.56,PURPLE)
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,9.86,6.26,1.18,.58,"FFFFFF",PURPLE,lw=1.1)
    add_text(slide,"知识增强回答",9.86,6.31,1.18,.44,10,INK,True,PP_ALIGN.CENTER)
    arrow(slide,11.05,6.56,11.55,6.56,PURPLE)
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,11.64,6.26,.96,.58,"E9F2E5",GREEN,lw=1.1)
    add_text(slide,"来源可追溯",11.64,6.31,.96,.44,10,GREEN,True,PP_ALIGN.CENTER)
    # index to online retrieval
    arrow(slide,11.45,5.29,5.05,6.29,GREEN,1.4)

    # Side vertical principle label, like the reference image.
    shape(slide,MSO_SHAPE.RECTANGLE,12.97,1.38,.34,4.18,GREEN_BG,GREEN,lw=1.1)
    t=add_text(slide,"敏感度引导的隐私—效用平衡",12.98,1.48,.30,3.98,12,GREEN,True,PP_ALIGN.CENTER)
    t.rotation=90

    prs.core_properties.title = "DP-RAG 中文版技术路线与原理图"
    prs.core_properties.subject = "参考论文图风格，全部原生形状可编辑"
    prs.core_properties.author = "Codex"
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
