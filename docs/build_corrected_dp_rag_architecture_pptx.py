"""Build the corrected, fully editable Chinese DP-RAG architecture figure."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from build_dp_rag_technical_route_pptx import (
    add_text, connector, dot, icon_document, icon_graph, icon_llm, icon_shield,
    rgb, shape,
)

OUT = Path(__file__).with_name("DP_RAG_自适应差分隐私检索框架_中文版可编辑.pptx")

INK="18202A"; MUTED="5E6B78"; WHITE="FFFFFF"; BG="FBFCFE"; LINE="CAD3DE"
BLUE="2166AC"; BLUE_L="EDF5FF"; GREEN="258246"; GREEN_L="EDF8F0"
PURPLE="6236A4"; PURPLE_L="F5F0FC"; GRAY="62666B"; GRAY_L="F3F4F5"
ORANGE="C85B19"; ORANGE_L="FFF3EA"; RED="A32929"; RED_L="FFF1F1"


def panel(slide,x,y,w,h,fill,line,title,title_color,icon=None):
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,fill,line,lw=1.15)
    connector(slide,x,y+.48,x+w,y+.48,line,1.0,arrow=False)
    if icon: icon(slide,x+.13,y+.08)
    add_text(slide,title,x+.12,y+.06,w-.24,.32,13,title_color,True,PP_ALIGN.CENTER)


def card(slide,x,y,w,h,fill=WHITE,line=LINE):
    return shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h,fill,line,lw=.85)


def arrow(slide,x1,y1,x2,y2,color=INK,width=1.25):
    return connector(slide,x1,y1,x2,y2,color,width,arrow=True)


def step_badge(slide,n,x,y,color=GREEN):
    shape(slide,MSO_SHAPE.OVAL,x,y,.24,.24,color,color,lw=0)
    add_text(slide,str(n),x,y,.24,.24,8,WHITE,True,PP_ALIGN.CENTER)


def tiles(slide,x,y,count=7,color=BLUE):
    colors=["D7E8FA","BBD4EF","E8D8CC","D8CDEC","C8DCF1"]
    for i in range(count):
        shape(slide,MSO_SHAPE.RECTANGLE,x+i*.18,y,.16,.19,colors[i%len(colors)],color,lw=.45)
    add_text(slide,"…",x+count*.18+.02,y-.02,.2,.2,9,INK,True,PP_ALIGN.CENTER)


def small_network(slide,x,y,color=BLUE):
    pts=[(x,y+.22),(x+.19,y),(x+.22,y+.42),(x+.43,y+.13),(x+.48,y+.41),(x+.68,y+.24)]
    for a,b in ((0,1),(0,2),(1,3),(2,3),(2,4),(3,5),(4,5)):
        connector(slide,*pts[a],*pts[b],color,.75,arrow=False)
    for px,py in pts: dot(slide,px-.035,py-.035,.08,WHITE,color)


def doc_icon(slide,x,y): icon_document(slide,x,y,BLUE)
def shield_green(slide,x,y): icon_shield(slide,x,y,GREEN)
def shield_purple(slide,x,y): icon_shield(slide,x,y,PURPLE)


def build():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=rgb(BG)
    add_text(slide,"面向隐私保护 RAG 的自适应差分隐私检索框架",1.35,.08,10.65,.38,20,INK,True,PP_ALIGN.CENTER)

    # Four principal columns.
    panel(slide,.12,.52,2.30,4.96,BLUE_L,BLUE,"数据拥有者（可信）",BLUE,doc_icon)
    panel(slide,2.55,.52,4.45,4.96,GREEN_L,GREEN,"自适应差分隐私编码器（核心方法）",GREEN,shield_green)
    panel(slide,7.13,.52,3.08,4.96,GRAY_L,GRAY,"不可信向量数据库",INK,None)
    panel(slide,10.34,.52,2.87,4.96,PURPLE_L,PURPLE,"用户与本地生成端",PURPLE,shield_purple)

    # Column 1: data preparation.
    add_text(slide,"原始文档",.86,1.13,.85,.24,11,INK,True,PP_ALIGN.CENTER)
    icon_document(slide,.93,1.40,BLUE)
    arrow(slide,1.27,1.93,1.27,2.18,BLUE)
    add_text(slide,"重叠分块",.86,2.17,.83,.23,10,INK,True,PP_ALIGN.CENTER)
    tiles(slide,.46,2.52,7,BLUE)
    arrow(slide,1.27,2.82,1.27,3.08,BLUE)
    small_network(slide,.91,3.16,BLUE)
    add_text(slide,"嵌入模型",.78,3.64,1.0,.22,10,INK,True,PP_ALIGN.CENTER)
    arrow(slide,1.27,3.91,1.27,4.14,BLUE)
    card(slide,.35,4.17,1.84,.74,WHITE,BLUE)
    add_text(slide,"原始语义向量",.49,4.28,1.56,.24,11,BLUE,True,PP_ALIGN.CENTER)
    tiles(slide,.56,4.59,7,BLUE)
    add_text(slide,"原文与生成模型均保持在本地",.42,5.03,1.70,.25,9,MUTED,False,PP_ALIGN.CENTER)
    arrow(slide,2.19,4.55,2.55,4.55,BLUE,1.6)

    # Column 2: corrected DP pipeline, six compact stages.
    stages=[
        ("混合隐私敏感度评估","规则匹配 + 关键词 + 语义判断","rᵢ ∈ [0.1,10] → sᵢ ∈ [0,1]"),
        ("动态隐私参数映射","εᵢ = f(sᵢ)，Δᵢ = g(sᵢ)","高敏感 → 小 εᵢ、大 Δᵢ"),
        ("加噪前 L2 范数裁剪","vᵢᶜ = vᵢ / max(1, ‖vᵢ‖₂ / C)","建立统一敏感度边界"),
        ("解析高斯噪声校准","求解 g(u*) = δ，σᵢ = u*Δᵢ","按分块敏感度自适应校准"),
        ("维度能量修正与加噪","σᵢ,dim = σᵢ·α/√d","zᵢ ~ N(0, σ²ᵢ,dim I)"),
        ("最终 L2 归一化","ṽᵢ = (vᵢᶜ + zᵢ)/‖vᵢᶜ + zᵢ‖₂","输出单位隐私向量"),
    ]
    y=.99
    for i,(title,formula,note) in enumerate(stages,1):
        h=.67
        card(slide,2.78,y,3.99,h,WHITE,"8FC59D")
        step_badge(slide,i,2.90,y+.10,GREEN)
        add_text(slide,title,3.22,y+.05,2.45,.22,10.5,GREEN,True,PP_ALIGN.CENTER)
        add_text(slide,formula,3.17,y+.28,2.63,.20,9.2,INK,False,PP_ALIGN.CENTER)
        add_text(slide,note,5.72,y+.12,.91,.36,8.1,MUTED,False,PP_ALIGN.CENTER)
        if i<len(stages): arrow(slide,4.77,y+h,4.77,y+h+.08,GREEN,1.0)
        y+=.73
    arrow(slide,7.00,4.55,7.13,4.55,GREEN,1.6)

    # Column 3: untrusted HNSW storage and retrieval.
    card(slide,7.37,1.04,2.60,.91,WHITE,"AAB0B5")
    add_text(slide,"存储：差分隐私保护向量",7.54,1.14,2.26,.23,10.5,INK,True,PP_ALIGN.CENTER)
    for i in range(26):
        x=7.77+(i%9)*.18; y=1.48+(i//9)*.13
        dot(slide,x,y,.055,"8F83D7" if i%3 else "75A5DF","8F83D7")
    arrow(slide,8.67,1.95,8.67,2.18,GRAY)
    card(slide,7.37,2.21,2.60,1.28,WHITE,"AAB0B5")
    add_text(slide,"HNSW 近似最近邻索引",7.55,2.31,2.24,.23,10.5,INK,True,PP_ALIGN.CENTER)
    icon_graph(slide,8.27,2.69)
    arrow(slide,8.67,3.49,8.67,3.72,GRAY)
    card(slide,7.37,3.75,2.60,1.22,WHITE,"AAB0B5")
    add_text(slide,"余弦相似度检索",7.61,3.85,2.12,.23,10.5,INK,True,PP_ALIGN.CENTER)
    add_text(slide,"返回 Top‑K 分块标识",7.59,4.20,2.16,.22,9.5,MUTED,False,PP_ALIGN.CENTER)
    for i,t in enumerate(("1","7","…","k")):
        shape(slide,MSO_SHAPE.RECTANGLE,7.74+i*.48,4.51,.48,.28,WHITE,"8D969E",lw=.65)
        add_text(slide,t,7.74+i*.48,4.51,.48,.28,9,INK,True,PP_ALIGN.CENTER)
    add_text(slide,"服务器可观察存储向量和查询，\n但差分隐私降低语义推断与重构风险",7.35,5.04,2.64,.32,8.5,RED,True,PP_ALIGN.CENTER)

    # Column 4: query stays clean; no noise on query.
    card(slide,10.57,1.03,2.41,.52,WHITE,"B99DE0")
    add_text(slide,"用户查询 q",10.72,1.13,2.10,.25,11,PURPLE,True,PP_ALIGN.CENTER)
    arrow(slide,11.77,1.55,11.77,1.76,PURPLE)
    card(slide,10.57,1.79,2.41,.72,WHITE,"B99DE0")
    small_network(slide,10.81,1.94,PURPLE)
    add_text(slide,"查询向量编码\n与 L2 归一化",11.52,1.91,1.25,.42,10,PURPLE,True,PP_ALIGN.CENTER)
    arrow(slide,11.77,2.51,11.77,2.70,PURPLE)
    card(slide,10.57,2.73,2.41,.55,WHITE,"B99DE0")
    add_text(slide,"检索请求（不对查询加噪）",10.71,2.83,2.12,.27,10,PURPLE,True,PP_ALIGN.CENTER)
    arrow(slide,10.57,3.00,10.21,3.00,PURPLE,1.4)
    arrow(slide,11.77,3.28,11.77,3.49,PURPLE)
    card(slide,10.57,3.52,2.41,.64,WHITE,"B99DE0")
    add_text(slide,"检索到的 Top‑K 文档分块",10.68,3.62,2.18,.25,10,PURPLE,True,PP_ALIGN.CENTER)
    tiles(slide,10.96,3.91,6,PURPLE)
    arrow(slide,11.77,4.16,11.77,4.36,ORANGE)
    card(slide,10.57,4.39,2.41,.46,ORANGE_L,"E7A06F")
    icon_llm(slide,10.70,4.36)
    add_text(slide,"本地 LLM 生成",11.50,4.46,1.22,.24,10,ORANGE,True,PP_ALIGN.CENTER)
    arrow(slide,11.77,4.85,11.77,5.02,ORANGE)
    card(slide,10.57,5.03,2.41,.32,WHITE,"B99DE0")
    add_text(slide,"知识增强回答",10.75,5.03,2.05,.32,10,PURPLE,True,PP_ALIGN.CENTER)

    # Bottom threat model and goals.
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,.12,5.65,13.09,1.70,RED_L,"E4A3A3",lw=.9)
    add_text(slide,"威胁模型、隐私边界与效用目标",4.46,5.69,4.42,.28,14,RED,True,PP_ALIGN.CENTER)
    connector(slide,.28,6.02,13.05,6.02,"E4A3A3",.7,arrow=False)
    # trust boundary
    add_text(slide,"信任边界",.42,6.11,1.22,.22,10,BLUE,True,PP_ALIGN.CENTER)
    connector(slide,.55,6.48,2.65,6.48,BLUE,1.0,True,False)
    add_text(slide,"可信：数据拥有者、本地 LLM",.42,6.57,2.44,.22,8.7,INK,False)
    connector(slide,.55,6.90,2.65,6.90,RED,1.0,True,False)
    add_text(slide,"不可信：外部向量数据库",.42,6.99,2.44,.22,8.7,INK,False)
    # adversary
    add_text(slide,"攻击者能力",3.15,6.11,1.42,.22,10,RED,True,PP_ALIGN.CENTER)
    add_text(slide,"• 知道嵌入模型与保护机制\n• 观察存储向量和查询向量\n• 尝试成员推断或语义重构",3.06,6.38,2.72,.72,8.6,INK,False)
    # guarantee
    add_text(slide,"隐私机制",6.18,6.11,1.42,.22,10,GREEN,True,PP_ALIGN.CENTER)
    add_text(slide,"按分块敏感度自适应校准 (εᵢ,δ) 参数，\n降低原始向量及文档语义的泄露风险。",5.98,6.39,3.08,.54,8.8,INK,False,PP_ALIGN.CENTER)
    card(slide,6.52,6.96,2.00,.27,WHITE,"D9AAAA")
    add_text(slide,"M(D) ≈ M(D′)",6.52,6.96,2.00,.27,10,INK,True,PP_ALIGN.CENTER)
    # utility
    add_text(slide,"效用与效率目标",9.47,6.11,1.60,.22,10,ORANGE,True,PP_ALIGN.CENTER)
    add_text(slide,"在隐私约束下保持 Recall / MRR，\n并利用 HNSW 提升查询与扩展效率。",9.13,6.39,3.47,.54,8.8,INK,False,PP_ALIGN.CENTER)
    card(slide,9.76,6.96,2.23,.27,WHITE,"D9AAAA")
    add_text(slide,"隐私保护 · 检索可用 · 可扩展",9.76,6.96,2.23,.27,9,RED,True,PP_ALIGN.CENTER)

    prs.core_properties.title="面向隐私保护 RAG 的自适应差分隐私检索框架"
    prs.core_properties.subject="中文版可编辑架构图，已按代码实现修正"
    prs.core_properties.author="Codex"
    prs.save(OUT); print(OUT)


if __name__=="__main__": build()
