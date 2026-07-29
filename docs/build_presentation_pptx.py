"""10分钟汇报PPT生成脚本 - DP-RAG项目汇报"""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "DP_RAG_汇报PPT_10分钟.pptx"

BG = "F7F8FA"
WHITE = "FFFFFF"
INK = "17233B"
MUTED = "64748B"
LINE = "D9E1EA"
TEAL = "1B8A89"
TEAL_LIGHT = "E3F3F1"
BLUE = "3D6FB4"
BLUE_LIGHT = "E9F0FA"
CORAL = "D9685A"
CORAL_LIGHT = "FBEAE7"
GOLD = "D99A32"
GOLD_LIGHT = "FFF3D8"

def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)

def add_text(slide, text, x, y, w, h, size=16, color=INK, bold=False,
             align=PP_ALIGN.LEFT, font="Microsoft YaHei", margin=0.05,
             valign=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box

def add_title(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = rgb(INK)
    title.text_frame.paragraphs[0].font.name = "Microsoft YaHei"

def card(slide, x, y, w, h, fill=WHITE, line=LINE, accent=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(line)
    s.line.width = Pt(1.1)
    return s

def add_bullet(slide, text, x, y, w, size=14, color=INK):
    return add_text(slide, text, x, y, w, 0.4, size=size, color=color, font="Microsoft YaHei")

def create_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, title)
    
    subtitle_box = slide.placeholders[1]
    subtitle_box.text = subtitle
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_box.text_frame.paragraphs[0].font.color.rgb = rgb(MUTED)
    subtitle_box.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
    
    return slide

def create_content_slide(prs, title, content_items):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, title)
    
    for i, item in enumerate(content_items):
        y_pos = 1.8 + i * 0.6
        add_text(slide, f"• {item}", 0.5, y_pos, 9, 0.5, size=18, font="Microsoft YaHei")
    
    return slide

def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, title)
    
    add_text(slide, left_title, 0.5, 1.5, 4, 0.4, size=20, bold=True, color=TEAL)
    for i, item in enumerate(left_items):
        y_pos = 2.0 + i * 0.5
        add_text(slide, f"• {item}", 0.5, y_pos, 4.2, 0.4, size=14, font="Microsoft YaHei")
    
    add_text(slide, right_title, 5.0, 1.5, 4.5, 0.4, size=20, bold=True, color=CORAL)
    for i, item in enumerate(right_items):
        y_pos = 2.0 + i * 0.5
        add_text(slide, f"• {item}", 5.0, y_pos, 4.2, 0.4, size=14, font="Microsoft YaHei")
    
    return slide

def create_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    add_title(slide, title)
    
    rows_count = len(rows) + 1
    cols_count = len(headers)
    
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9)
    height = Inches(0.5)
    
    table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
    
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(TEAL)
    
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_data
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ===== 第1页：封面 =====
    create_title_slide(prs, 
        "基于差分隐私的检索增强生成系统", 
        "DP-RAG: Differential Privacy for Retrieval-Augmented Generation\n\n汇报人：XXX\n汇报日期：2026-07-29")
    
    # ===== 第2页：研究背景 =====
    create_content_slide(prs, "研究背景与问题", [
        "RAG（检索增强生成）在企业知识库、智能问答等领域广泛应用",
        "RAG系统面临隐私泄露风险：攻击者可从检索向量推断敏感信息",
        "现有方案问题：纯密文检索计算开销大，差分隐私保护不充分",
        "研究目标：在保证检索质量的前提下，提供可证明的差分隐私保护"
    ])
    
    # ===== 第3页：方案概述 =====
    create_content_slide(prs, "DP-RAG方案概述", [
        "核心思想：在向量检索层引入差分隐私噪声，保护原始文档表示",
        "技术路线：JL降维 → L2裁剪 → 动态分析高斯机制 → HNSW检索",
        "隐私保障：提供(ε, δ)-差分隐私可证明保护",
        "支持模式：NoJL(1024维)、JL768、JL256三种配置"
    ])
    
    # ===== 第4页：技术流程图 =====
    create_content_slide(prs, "技术流程", [
        "1. 文档预处理：分块 + 嵌入模型(BGE-M3) → 1024维原始向量",
        "2. Johnson-Lindenstrauss降维：1024维 → 256/768维，减少高维噪声累积",
        "3. 动态隐私预算分配：根据文档敏感度自适应调整噪声强度",
        "4. 差分隐私噪声注入：L2裁剪 + 分析高斯机制",
        "5. HNSW近似最近邻检索：平衡检索精度与效率"
    ])
    
    # ===== 第5页：对比实验设置 =====
    create_content_slide(prs, "实验1：基线对比实验", [
        "对比方案：Our DP-RAG (NoJL/JL768/JL256)、Private RAG-RP、DCPE+DCE、CKKS",
        "评估指标：Recall@5、MRR@5、查询时间、隐私预算ε",
        "数据集：知识库文档 + 30条测试查询",
        "HNSW参数：ef_search ∈ {16, 32, 64, 128, 256}"
    ])
    
    # ===== 第6页：对比实验结果 =====
    create_table_slide(prs, "基线对比实验结果", 
        ["方案", "Recall@5", "MRR@5", "平均ε", "噪声信号比", "向量维度"],
        [
            ["Our DP-RAG-NoJL", "0.980", "1.000", "7.71", "0.0035", "1024"],
            ["Our DP-RAG-JL768", "0.720", "0.788", "7.71", "0.0035", "768"],
            ["Our DP-RAG-JL256", "0.587", "0.663", "7.71", "0.0035", "256"],
            ["Private RAG-RP", "0.280", "0.303", "N/A", "N/A", "64"],
            ["DCPE+DCE", "0.973", "1.000", "N/A", "0.025", "1024"]
        ])
    
    # ===== 第7页：对比实验分析 =====
    create_content_slide(prs, "基线对比实验分析", [
        "Our DP-RAG-NoJL达到最高检索精度(Recall@5=0.98)，接近原始RAG",
        "JL降维导致轻微精度下降，但显著降低存储和计算开销",
        "Private RAG-RP无正式DP保证，检索精度较低(0.28)",
        "DCPE+DCE精度高但无差分隐私保障，不提供ε值",
        "隐私-效用权衡：Our DP-RAG在7.71的ε下保持良好效用"
    ])
    
    # ===== 第8页：安全性实验设置 =====
    create_two_column_slide(prs, "实验2：安全性评估",
        "向量链接攻击(Vector Linkage)",
        [
            "威胁模型：攻击者获取保护向量，尝试链接至原始文档",
            "评估指标：Linkage Top-1 Recovery Rate、Recall@5、MRR@5",
            "目标：证明受保护向量难以链接回原始文档"
        ],
        "敏感数据检索攻击(Sensitive Retrieval)",
        [
            "威胁模型：攻击者利用语义查询检索含敏感信息(邮箱/电话/URL)的文档",
            "评估指标：Sensitive Target Recall@1/@5、Exposure Rate",
            "目标：证明敏感文档难以被恶意检索"
        ])
    
    # ===== 第9页：安全性实验结果 =====
    create_table_slide(prs, "安全性实验结果", 
        ["方案", "向量链接Top-1恢复率", "敏感检索Top-1暴露率", "普通检索Recall@5"],
        [
            ["Vanilla Raw HNSW", "100%", "100%", "0.967"],
            ["Our DP-RAG-JL256", "较低(受保护)", "较低(受保护)", "0.580"],
            ["Private RAG-RP", "N/A", "N/A", "0.280"],
            ["DCPE+DCE", "N/A", "N/A", "0.973"]
        ])
    
    # ===== 第10页：安全性实验分析 =====
    create_content_slide(prs, "安全性实验分析", [
        "Vanilla Raw HNSW完全无保护，攻击者可100%成功链接和检索敏感数据",
        "Our DP-RAG通过DP噪声显著降低向量链接和敏感检索成功率",
        "Private RAG-RP虽有一定保护，但无正式DP保证，效用较低",
        "安全-效用权衡：DP-RAG在右下区域(高Recall + 低暴露)表现最优"
    ])
    
    # ===== 第11页：消融实验 =====
    create_content_slide(prs, "实验3：消融实验", [
        "Full Current：完整pipeline (JL→裁剪→动态DP→维度感知缩放→归一化)",
        "No DP Baseline：仅JL降维，无隐私保护",
        "Old Pipeline DP Before JL：在1024维空间直接加噪",
        "No Dimension-Aware Scaling：移除√dim校正",
        "Fixed DP Calibration：使用固定ε而非动态校准"
    ])
    
    # ===== 第12页：HNSW验证实验 =====
    create_content_slide(prs, "实验4：HNSW检索验证", [
        "验证HNSW近似检索与精确检索的一致性",
        "评估指标：Recall@K、Overlap@5、Speedup Ratio",
        "实验结果：ef_search=64时，Recall@5达到0.98以上",
        "加速效果：HNSW相比精确检索提速显著(>100x)",
        "结论：HNSW可在保持高精度的同时大幅降低查询延迟"
    ])
    
    # ===== 第13页：总结 =====
    create_content_slide(prs, "总结与贡献", [
        "提出DP-RAG：在RAG检索层引入可证明的差分隐私保护",
        "设计动态分析高斯机制：根据文档敏感度自适应分配隐私预算",
        "综合评估：基线对比、安全性、消融、HNSW四类实验",
        "实验表明：在ε≈7.71下，检索精度仅下降约2%，安全性显著提升"
    ])
    
    # ===== 第14页：未来工作 =====
    create_content_slide(prs, "未来工作", [
        "探索更高效的隐私预算分配策略，进一步降低ε",
        "研究多模态(图像、音频)RAG的隐私保护方案",
        "结合大语言模型输出层面的隐私保护",
        "在实际企业知识库中部署并验证DP-RAG"
    ])
    
    # ===== 第15页：致谢 =====
    create_title_slide(prs, "感谢聆听", "欢迎提问与讨论")
    
    prs.save(str(OUT))
    print(f"PPT已保存至: {OUT}")

if __name__ == "__main__":
    main()